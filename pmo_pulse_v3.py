# -*- coding: utf-8 -*-
"""
PMO Pulse Narrative Driven App (v1.0 - Part 1)

Refactored version of pmo_pulse_v2.py focusing on narrative flow,
visual appeal, Arcadis branding, and sidebar configuration.
This version focuses on Earned Value Management (EVM) KPIs.

Based on the framework developed for RiskLens Pro and CostBench.

Requires: pip install streamlit pandas numpy plotly python-dateutil joblib openpyxl scipy python-pptx fpdf2 (Optional for PDF)
"""

# ==============================================================================
# Part 1: Imports, Configuration, Helpers, Data Loading, State Init, Sidebar, Welcome Tab
# ==============================================================================

# --- Import Libraries ---
import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import datetime
from dateutil.relativedelta import relativedelta
import time
import logging
import io
import os
# from fpdf import FPDF # Optional for PDF export

# --- Basic Logging Configuration ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# --- Constants & Configuration ---

# Arcadis Branding (Consistent with previous examples)
ARCADIS_ORANGE = "#E67300"
ARCADIS_BLACK = "#000000"
ARCADIS_GREY = "#6c757d" # Slightly darker grey for text
ARCADIS_DARK_GREY = "#646469" # From RiskLens for consistency
ARCADIS_WHITE = "#FFFFFF"
ARCADIS_LIGHT_GREY = "#F5F5F5" # Consistent light grey background
ARCADIS_SECONDARY_PALETTE = ["#00A3A1", ARCADIS_DARK_GREY, "#D6D6D8"] # Teal, Dark Grey, Mid Grey
COLOR_SUCCESS = "#2ECC71"
COLOR_WARNING = "#F1C40F"
COLOR_DANGER = "#E74C3C"
COLOR_INFO = "#3498DB"

# Placeholder Logos (Replace with actual URLs)
ARCADIS_LOGO_URL_WIDE = f"https://placehold.co/200x50/{ARCADIS_WHITE[1:]}/{ARCADIS_ORANGE[1:]}?text=Arcadis+Logo"
ARCADIS_LOGO_ICON_URL = f"https://placehold.co/50x50/{ARCADIS_WHITE[1:]}/{ARCADIS_ORANGE[1:]}?text=A"
ARCADIS_FAVICON = ARCADIS_LOGO_ICON_URL

# Plotly Template
PLOTLY_TEMPLATE = "plotly_white" # Use Plotly's white template as a base

# --- Page Configuration ---
# IMPORTANT: Must be the first Streamlit command
st.set_page_config(
    page_title="Arcadis | PMO Pulse",
    page_icon=ARCADIS_FAVICON,
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- Load Custom CSS ---
# Adapted CSS for PMO Pulse narrative style
st.markdown(f"""
<style>
    /* Base font */
    html, body, [class*="st-"] {{ font-family: 'Arial', sans-serif; }}

    /* Button Styling */
    .stButton>button {{
        border-radius: 8px; border: 1px solid {ARCADIS_ORANGE}; background-color: {ARCADIS_ORANGE}; color: white;
        transition: background-color 0.3s ease, border-color 0.3s ease; font-weight: bold; padding: 0.5rem 1rem;
        width: auto; display: inline-block; margin-right: 10px; margin-top: 10px;
    }}
    .stButton>button:hover {{ background-color: #D06300; border-color: #D06300; color: white; }}
    .stButton>button:active {{ background-color: #B85A00; border-color: #B85A00; }}
    .stButton>button:focus {{ outline: none; box-shadow: 0 0 0 2px rgba(230, 115, 0, 0.5); }}
    .stButton>button:disabled {{ background-color: #cccccc; color: #666666; border-color: #cccccc; }}

    /* Sidebar Button Styling */
    [data-testid="stSidebar"] .stButton>button {{ width: 95%; margin-bottom: 10px; }}

    /* Download Button Styling */
    .stDownloadButton>button {{
        background-color: {ARCADIS_DARK_GREY}; border-color: {ARCADIS_DARK_GREY}; color: white; font-weight: bold;
        padding: 0.5rem 1rem; border-radius: 8px; width: auto; display: inline-block; margin-right: 10px; margin-bottom: 10px;
    }}
    .stDownloadButton>button:hover {{ background-color: #505055; border-color: #505055; }}

    /* Sidebar Styling */
    .stSidebar {{ background-color: {ARCADIS_WHITE}; border-right: 1px solid #D6D6D8; }}
    [data-testid="stSidebarNav"] {{ padding-top: 0rem; }}
    [data-testid="stSidebarUserContent"] {{ padding-top: 1rem; }}

    /* Headings */
    h1, h2 {{ color: {ARCADIS_BLACK}; font-weight: bold; }}
    .main h3 {{ /* Custom H3 for main content sections */
        color: {ARCADIS_ORANGE}; font-weight: bold; border-bottom: 2px solid {ARCADIS_ORANGE};
        padding-bottom: 5px; margin-top: 1.5rem; margin-bottom: 1rem;
    }}
    h4, h5, h6 {{ color: {ARCADIS_DARK_GREY}; }}

    /* Metric Styling (Tiles) */
    .stMetric {{
        background-color: {ARCADIS_WHITE}; border: 1px solid #D6D6D8; border-left: 5px solid {ARCADIS_ORANGE};
        border-radius: 8px; padding: 15px 20px; box-shadow: 0 2px 4px rgba(0,0,0,0.05); margin-bottom: 15px;
    }}
    .stMetric > label {{ font-weight: bold; color: {ARCADIS_DARK_GREY}; font-size: 0.95rem; }}
    .stMetric > div[data-testid="stMetricValue"] {{ font-size: 2em; font-weight: bold; color: {ARCADIS_BLACK}; }}
    .stMetric > div[data-testid="stMetricDelta"] {{ font-size: 0.9em; }}
    /* Ensure delta colors are correct */
    .stMetric .st-emotion-cache-1g6goys {{ color: green !important; }} /* Default positive delta */
    .stMetric .st-emotion-cache-1g6goys.st-emotion-cache-1g6goys-red {{ color: red !important; }} /* Negative delta */


    /* Tab Styling */
    .stTabs [data-baseweb="tab-list"] button[aria-selected="true"] {{
        background-color: {ARCADIS_ORANGE}; color: white; border-radius: 8px 8px 0 0; font-weight: bold; border-bottom: none;
    }}
    .stTabs [data-baseweb="tab-list"] button {{
        border-radius: 8px 8px 0 0; color: {ARCADIS_DARK_GREY}; background-color: #E0E0E0; border-bottom: none;
    }}
    .stTabs [data-baseweb="tab-list"] {{ border-bottom: 2px solid {ARCADIS_ORANGE}; padding-bottom: 0; }}
    .stTabs [data-baseweb="tab-panel"] {{ background-color: {ARCADIS_LIGHT_GREY}; padding-top: 25px; border: none; }}

    /* Containers */
    .stVerticalBlock {{ padding-bottom: 1rem; }}
    /* Style containers used for sections */
    div[data-testid="stVerticalBlock"]>div[style*="flex-direction: column;"]>div[data-testid="stVerticalBlock"],
    div[data-testid="stVerticalBlock"]>div[style*="flex-direction: column;"]>div[data-testid="stHorizontalBlock"] {{
        border-radius: 8px !important; border: 1px solid #D6D6D8 !important; padding: 20px !important;
        margin-bottom: 20px !important; background-color: {ARCADIS_WHITE} !important; box-shadow: 0 2px 4px rgba(0,0,0,0.05) !important;
    }}
    /* Welcome page specific container */
    .welcome-section {{
        background-color: {ARCADIS_WHITE}; border: 1px solid #D6D6D8; border-radius: 8px;
        padding: 25px; margin-bottom: 20px; box-shadow: 0 2px 4px rgba(0,0,0,0.05);
    }}
     .welcome-section h3 {{ border: none; margin-bottom: 10px; margin-top: 0.5rem; }}

    /* Arcadis Logo in Sidebar */
    [data-testid="stSidebarNav"]::before {{
        content: ""; display: block; background-image: url({ARCADIS_LOGO_URL_WIDE});
        background-size: contain; background-repeat: no-repeat; background-position: center 10px;
        height: 60px; margin-bottom: 10px;
    }}

    /* Expander in Sidebar */
    [data-testid="stSidebar"] .stExpander {{
        border: none !important; border-radius: 0px !important; background-color: transparent !important; margin-bottom: 0px;
        border-top: 1px solid #eee !important;
    }}
    [data-testid="stSidebar"] .stExpander header {{
        font-weight: bold; color: {ARCADIS_BLACK}; background-color: transparent !important; border-radius: 0 !important;
        padding: 10px 0px !important;
    }}
     [data-testid="stSidebar"] .stExpander div[data-testid="stExpanderDetails"] {{
         padding-left: 10px !important;
     }}

    /* Dataframes */
    .stDataFrame {{ border-radius: 8px; overflow: hidden; border: 1px solid #e0e0e0; }}
    .stDataFrame thead th {{ background-color: {ARCADIS_DARK_GREY}; color: white; font-weight: bold; }}
    .stDataFrame td, .stDataFrame th {{ border-bottom: 1px solid #e0e0e0; border-right: none; }}

    /* Markdown links */
    a {{ color: {ARCADIS_ORANGE}; }}

    /* Main area padding */
    .main .block-container {{ padding: 2rem; }}

    /* Lists for capabilities/questions */
    .styled-list li {{ margin-bottom: 10px; line-height: 1.6; color: {ARCADIS_DARK_GREY}; }}
    .styled-list li b {{ color: {ARCADIS_BLACK}; }}
    .styled-list li i {{ color: {ARCADIS_ORANGE}; font-style: normal; font-weight: bold; }} /* Tab names */

    /* Callout boxes for dashboard alerts */
    .callout {{
        padding: 15px; margin-bottom: 15px; border: 1px solid transparent; border-radius: 8px;
        display: flex; align-items: center; border-left-width: 5px;
    }}
    .callout h4 {{ margin-top: 0; margin-bottom: 5px; color: inherit; border: none; font-size: 1.1em; }}
    .callout p {{ margin-bottom: 0; font-size: 0.9em; line-height: 1.4; }}
    .callout .icon {{ font-size: 24px; margin-right: 15px; }}
    .callout-danger {{ background-color: #f8d7da; border-color: #f5c6cb; color: #721c24; }}
    .callout-danger .icon {{ color: #721c24; }}
    .callout-warning {{ background-color: #fff3cd; border-color: #ffeeba; color: #856404; }}
    .callout-warning .icon {{ color: #856404; }}
    .callout-success {{ background-color: #d4edda; border-color: #c3e6cb; color: #155724; }}
    .callout-success .icon {{ color: #155724; }}

    /* Attention list items */
     .attention-list-item {{
        padding: 8px 12px; margin-bottom: 5px; border-radius: 4px; border: 1px solid #eee;
        display: flex; justify-content: space-between; align-items: center; font-size: 0.9em;
    }}
    .attention-list-item span {{ font-size: 0.85em; color: #6c757d; margin-left: 10px; }}
    .attention-list-item.spi {{ border-left: 3px solid {ARCADIS_SECONDARY_PALETTE[0]}; }} /* Teal */
    .attention-list-item.cpi {{ border-left: 3px solid {ARCADIS_ORANGE}; }}
    .attention-list-item.risk {{ border-left: 3px solid {COLOR_DANGER}; }}

</style>
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.0.0/css/all.min.css">
""", unsafe_allow_html=True)


# --- Helper Functions ---

# Placeholder Security Functions
def authenticate_user():
    """Placeholder for user authentication."""
    st.session_state['authenticated'] = True
    st.session_state['user_role'] = "viewer" # Default role
    return True

def authorize_user(required_role="viewer"):
    """Placeholder for role-based authorization."""
    if not st.session_state.get('authenticated', False): return False
    user_role = st.session_state.get('user_role', 'viewer')
    role_levels = {"admin": 3, "editor": 2, "viewer": 1}
    required_level = role_levels.get(required_role, 1)
    user_level = role_levels.get(user_role, 0)
    return user_level >= required_level

# Placeholder Data Connection Functions
def connect_to_database():
    logging.info("Attempting to connect to database (Placeholder)...")
    return "dummy_connection" # Simulate success

def fetch_from_api(api_endpoint):
    logging.info(f"Attempting to fetch data from API: {api_endpoint} (Placeholder)...")
    return {"data": []} # Simulate success

# Data Validation
def validate_data(df, df_name):
    """Performs basic data validation checks."""
    logging.info(f"Performing data validation for {df_name}...")
    if not isinstance(df, pd.DataFrame):
         logging.error(f"Validation Error: {df_name} is not a DataFrame.")
         return False, f"{df_name} is not a DataFrame."

    if df.empty:
        if df_name in ['changes']: # Optional dataframes can be empty
             logging.info(f"{df_name} dataframe is empty, which is acceptable.")
             return True, ""
        logging.warning(f"{df_name} dataframe is empty.")
        if df_name in ['projects', 'tasks']: # Essential dataframes should not be empty
             return False, f"{df_name} is empty (Essential Data Missing)."
        else:
             return True, ""

    required_columns = {
        'projects': ['project_id', 'project_name', 'sector', 'budget', 'start_date', 'planned_end_date', 'status'],
        'tasks': ['task_id', 'project_id', 'planned_start', 'planned_end', 'planned_cost', 'actual_cost', 'earned_value'],
        'risks': ['risk_id', 'project_id', 'probability', 'impact_cost'],
        'changes': ['change_id', 'project_id', 'status', 'impact_cost', 'impact_schedule_days'],
        'history': ['project_id', 'month', 'cpi', 'spi'],
        'benefits': ['Month', 'ReportingTimeSaved_hrs', 'CostOverrunsAvoided_k', 'ForecastAccuracy_perc']
    }

    if df_name not in required_columns:
        logging.warning(f"No validation rules defined for {df_name}.")
        return True, ""

    missing_cols = [col for col in required_columns[df_name] if col not in df.columns]
    if missing_cols:
        msg = f"Validation Error in {df_name}: Missing required columns: {', '.join(missing_cols)}"
        logging.error(msg)
        if df_name in ['changes']:
             logging.warning(msg + " (Non-fatal for optional data)")
             return True, "" # Treat as valid but log warning
        return False, msg

    # Add more specific checks (data types, ranges, consistency) here if needed
    # Example: Check if budget is numeric in projects
    if df_name == 'projects':
        if 'budget' in df.columns and not pd.api.types.is_numeric_dtype(df['budget']):
            msg = f"Validation Error in {df_name}: 'budget' column is not numeric."
            logging.error(msg); return False, msg
    # Example: Check date consistency
    if df_name == 'projects':
         if 'start_date' in df.columns and 'planned_end_date' in df.columns:
              start_dates = pd.to_datetime(df['start_date'], errors='coerce')
              end_dates = pd.to_datetime(df['planned_end_date'], errors='coerce')
              if (end_dates < start_dates).any():
                   logging.warning(f"Validation Warning in {df_name}: Some planned_end_date are before start_date.")

    logging.info(f"Basic validation passed for {df_name}.")
    return True, ""

# Mock Data Generation
@st.cache_data # Cache the mock data generation
def generate_mock_data(num_projects=40, tasks_per_project=20, risks_per_project=7, history_months=12):
    """Generates mock project data (Adapted from original script)."""
    logging.info("Generating mock data...")
    np.random.seed(42)
    today = datetime.date.today()
    first_possible_start = today - relativedelta(years=3)
    project_ids = np.arange(1, num_projects + 1)
    sectors = np.random.choice(['Infrastructure', 'Buildings', 'Water', 'Environment', 'Energy Transition', 'Digital Transformation'], num_projects)
    statuses = np.random.choice(['On Track', 'Minor Issues', 'At Risk', 'Delayed', 'Completed'], num_projects, p=[0.35, 0.2, 0.15, 0.15, 0.15])
    budgets = np.random.randint(500_000, 30_000_000, num_projects)
    start_dates = [first_possible_start + datetime.timedelta(days=np.random.randint(0, max(1, (today - first_possible_start).days - 180))) for _ in range(num_projects)]
    planned_durations = np.random.randint(180, 1200, num_projects)
    end_dates = [sd + datetime.timedelta(days=int(pd)) for sd, pd in zip(start_dates, planned_durations)]
    target_cpi = np.random.uniform(0.95, 1.02, num_projects)
    target_spi = np.random.uniform(0.98, 1.02, num_projects)
    projects_df = pd.DataFrame({
        'project_id': project_ids,
        'project_name': [f'Project {chr(65+(i%26))}{i//26 if i//26 > 0 else ""} ({sectors[i][0:3]})' for i in range(num_projects)],
        'sector': sectors, 'budget': budgets, 'start_date': pd.to_datetime(start_dates),
        'planned_end_date': pd.to_datetime(end_dates), 'status': statuses,
        'project_manager': np.random.choice(['Alice Smith', 'Bob Johnson', 'Charlie Brown', 'Diana Prince', 'Ethan Hunt', 'Fiona Glenanne'], num_projects),
        'planned_duration_days': planned_durations, 'target_cpi': target_cpi, 'target_spi': target_spi,
        'strategic_alignment': np.random.choice(['High', 'Medium', 'Low'], num_projects, p=[0.5, 0.4, 0.1]),
        'collaboration_link': [f"https://example.com/projects/{pid}" for pid in project_ids]
    })
    historical_data = []
    for pid in project_ids:
        proj = projects_df[projects_df['project_id'] == pid].iloc[0]
        base_cpi = np.random.normal(loc=1.0, scale=0.1); base_spi = np.random.normal(loc=1.0, scale=0.08)
        if proj['status'] == 'Delayed': base_spi -= 0.15; base_cpi -= 0.1
        elif proj['status'] == 'At Risk': base_spi -= 0.08; base_cpi -= 0.05
        elif proj['status'] == 'Minor Issues': base_spi -= 0.03; base_cpi -= 0.02
        for i in range(history_months):
            month_date = today - relativedelta(months=i)
            cpi = max(0.5, min(1.5, base_cpi + np.random.normal(scale=0.03) - (i * 0.002)))
            spi = max(0.5, min(1.5, base_spi + np.random.normal(scale=0.02) - (i * 0.003)))
            if pd.notna(proj['start_date']) and month_date >= proj['start_date'].date():
                historical_data.append({'project_id': pid, 'month': month_date.strftime('%Y-%m'), 'cpi': cpi, 'spi': spi})
    history_df = pd.DataFrame(historical_data)
    all_tasks_data = []
    task_id_counter = 1
    for pid in project_ids:
        proj = projects_df[projects_df['project_id'] == pid].iloc[0]
        proj_start = proj['start_date']; proj_end = proj['planned_end_date']
        if pd.isna(proj_start) or pd.isna(proj_end) or proj_end <= proj_start: proj_duration = 0
        else: proj_duration = (proj_end - proj_start).days
        proj_duration_offset_base = max(1, proj_duration - 30)
        for i in range(tasks_per_project):
            task_planned_start_offset = np.random.randint(0, max(1, proj_duration_offset_base))
            task_planned_duration = np.random.randint(10, max(11, proj_duration // tasks_per_project if tasks_per_project > 0 else 90))
            task_planned_start = proj_start + datetime.timedelta(days=task_planned_start_offset)
            task_planned_end = task_planned_start + datetime.timedelta(days=task_planned_duration)
            delay_factor = 0; cost_factor = 1.0; completion_factor = np.random.uniform(0.8, 1.0)
            if proj['status'] == 'Minor Issues': delay_factor = np.random.randint(0, 7); cost_factor = np.random.uniform(1.0, 1.07); completion_factor = np.random.uniform(0.7, 0.95)
            elif proj['status'] == 'At Risk': delay_factor = np.random.randint(2, 15); cost_factor = np.random.uniform(1.03, 1.18); completion_factor = np.random.uniform(0.5, 0.85)
            elif proj['status'] == 'Delayed': delay_factor = np.random.randint(10, 30); cost_factor = np.random.uniform(1.08, 1.30); completion_factor = np.random.uniform(0.3, 0.7)
            elif proj['status'] == 'Completed': delay_factor = np.random.randint(-5, 7); cost_factor = np.random.uniform(0.95, 1.1); completion_factor = 1.0
            task_actual_start = pd.NaT
            if pd.notna(task_planned_start) and task_planned_start.date() <= today and np.random.rand() > 0.05 :
                 task_actual_start = task_planned_start + datetime.timedelta(days=np.random.randint(0, max(1, delay_factor // 2 + 1)))
            task_actual_end = pd.NaT
            if pd.notna(task_actual_start) and completion_factor == 1.0:
                task_actual_end = task_planned_end + datetime.timedelta(days=delay_factor + np.random.randint(-2, 3))
                if pd.notna(task_actual_end) and task_actual_end < task_actual_start:
                    task_actual_end = task_actual_start + datetime.timedelta(days=np.random.randint(1, task_planned_duration))
            future_limit = pd.Timestamp(today + datetime.timedelta(days=730))
            task_actual_end = pd.Timestamp(min(task_actual_end, future_limit)) if pd.notna(task_actual_end) else pd.NaT
            task_planned_end = pd.Timestamp(min(task_planned_end, future_limit)) if pd.notna(task_planned_end) else pd.NaT
            if proj_duration > 0 and tasks_per_project > 0: task_planned_cost = max(1000, proj['budget'] * (task_planned_duration / proj_duration) * np.random.uniform(0.6, 1.4) / tasks_per_project)
            else: task_planned_cost = max(1000, proj['budget'] * np.random.uniform(0.01, 0.05))
            task_actual_cost = 0; task_earned_value = 0; percent_complete = 0
            if pd.notna(task_actual_start):
                days_since_start = (today - task_actual_start.date()).days if task_actual_start.date() <= today else 0
                planned_days = (task_planned_end.date() - task_planned_start.date()).days if pd.notna(task_planned_end) and pd.notna(task_planned_start) else 0
                if planned_days > 0 and proj['status'] != 'Completed':
                     time_progress = min(1.0, max(0.0, days_since_start / planned_days))
                     if proj['status'] == 'Delayed': time_progress *= 0.7
                     elif proj['status'] == 'At Risk': time_progress *= 0.85
                     elif proj['status'] == 'Minor Issues': time_progress *= 0.95
                     completion_factor = np.clip(time_progress + np.random.normal(scale=0.1), 0, 1.0)
                elif proj['status'] == 'Completed': completion_factor = 1.0
                task_actual_cost = task_planned_cost * cost_factor * np.random.uniform(0.85, 1.15) * completion_factor
                task_earned_value = task_planned_cost * completion_factor
                percent_complete = completion_factor
            all_tasks_data.append({
                'task_id': task_id_counter, 'project_id': pid, 'task_name': f'Task {i+1}',
                'planned_start': task_planned_start, 'planned_end': task_planned_end,
                'actual_start': task_actual_start,
                'actual_end': task_actual_end if percent_complete == 1.0 else pd.NaT,
                'planned_cost': task_planned_cost, 'actual_cost': task_actual_cost,
                'earned_value': task_earned_value, 'percent_complete': percent_complete,
                'resource': np.random.choice(['Team Alpha', 'Team Bravo', 'Contractor Zeta', 'Internal Experts', 'Specialist Gamma', 'Shared Services'], 1)[0]
            })
            task_id_counter += 1
    tasks_df = pd.DataFrame(all_tasks_data)
    tasks_df[['planned_start', 'planned_end', 'actual_start', 'actual_end']] = tasks_df[['planned_start', 'planned_end', 'actual_start', 'actual_end']].apply(pd.to_datetime, errors='coerce')
    all_risks_data = []
    risk_id_counter = 1
    for pid in project_ids:
        proj = projects_df[projects_df['project_id'] == pid].iloc[0]
        for i in range(risks_per_project):
            prob = np.random.beta(2, 5)
            min_impact = max(1000, int(proj['budget']*0.005)); max_impact = max(min_impact + 1000, int(proj['budget']*0.20))
            impact_cost = np.random.randint(min_impact, max_impact)
            risk_status = np.random.choice(['Open', 'Mitigating', 'Closed', 'Realized'], p=[0.5, 0.25, 0.15, 0.1])
            all_risks_data.append({
                'risk_id': risk_id_counter, 'project_id': pid,
                'risk_description': f'Risk {i+1} - {np.random.choice(["Scope Creep", "Resource Availability", "Technical Challenge", "Supplier Delay", "Regulatory Change", "Stakeholder Dissatisfaction", "Funding Cut"])}',
                'probability': prob, 'impact_cost': impact_cost, 'risk_score': prob * impact_cost,
                'mitigation_plan': f'Mitigation Strategy {i+1}' if np.random.rand() > 0.2 else 'Awaiting Plan',
                'risk_status': risk_status, 'owner': np.random.choice(['PM', 'Tech Lead', 'Client', 'Finance'], 1)[0],
                'collaboration_link': f"https://example.com/risks/{risk_id_counter}"
            })
            risk_id_counter += 1
    risks_df = pd.DataFrame(all_risks_data)
    all_changes_data = []
    change_id_counter = 1
    for pid in project_ids:
        num_changes = np.random.randint(0, 8)
        proj = projects_df[projects_df['project_id'] == pid].iloc[0]
        if pd.notna(proj['start_date']) and pd.notna(proj['planned_end_date']) and proj['planned_end_date'] > proj['start_date']:
            date_range_days = (proj['planned_end_date'] - proj['start_date']).days
            min_submit_offset = 20; max_submit_offset = max(min_submit_offset + 1, date_range_days - 20)
            if max_submit_offset > min_submit_offset:
                for i in range(num_changes):
                    min_change_impact = max(100, int(proj['budget']*0.001)); max_change_impact = max(min_change_impact + 100, int(proj['budget']*0.08))
                    impact_cost_change = np.random.randint(min_change_impact, max_change_impact)
                    submit_date = proj['start_date'] + datetime.timedelta(days=np.random.randint(min_submit_offset, max_submit_offset))
                    submit_date = min(submit_date.date(), today)
                    all_changes_data.append({
                        'change_id': change_id_counter, 'project_id': pid, 'description': f'Change Request {i+1}',
                        'impact_cost': impact_cost_change, 'impact_schedule_days': np.random.randint(-5, 45),
                        'status': np.random.choice(['Submitted', 'Approved', 'Rejected', 'Implemented', 'On Hold'], p=[0.3, 0.4, 0.1, 0.15, 0.05]),
                        'date_submitted': pd.to_datetime(submit_date),
                         'collaboration_link': f"https://example.com/changes/{change_id_counter}"
                    })
                    change_id_counter += 1
    changes_df = pd.DataFrame(all_changes_data)
    if not changes_df.empty: changes_df['date_submitted'] = pd.to_datetime(changes_df['date_submitted'], errors='coerce')
    benefits_data = {
        'Month': [(today - relativedelta(months=i)).strftime('%Y-%m') for i in range(history_months)][::-1],
        'ReportingTimeSaved_hrs': (np.linspace(5, 40, history_months) + np.random.normal(0, 3, history_months)).clip(min=0),
        'CostOverrunsAvoided_k': (np.cumsum(np.random.uniform(10, 50, history_months)) + np.random.normal(0, 20, history_months)).clip(min=0),
        'ForecastAccuracy_perc': np.clip(np.linspace(60, 85, history_months) + np.random.normal(0, 4, history_months), 50, 95)
    }
    benefits_df = pd.DataFrame(benefits_data)
    logging.info("Mock data generation complete.")
    return projects_df, tasks_df, risks_df, changes_df, history_df, benefits_df

# Data Loading Wrapper
@st.cache_data(ttl=3600) # Cache data for 1 hour
def load_data(use_mock_data=True):
    """Loads data from specified sources or generates mock data."""
    logging.info(f"Attempting to load data. Use Mock Data: {use_mock_data}")
    data_loaded_successfully = False
    validation_errors = []
    dfs = {}
    empty_dfs = { # Define structure for empty DFs
        'projects': pd.DataFrame(columns=['project_id', 'project_name', 'sector', 'budget', 'start_date', 'planned_end_date', 'status', 'project_manager', 'planned_duration_days', 'target_cpi', 'target_spi', 'strategic_alignment', 'collaboration_link']),
        'tasks': pd.DataFrame(columns=['task_id', 'project_id', 'task_name', 'planned_start', 'planned_end', 'actual_start', 'actual_end', 'planned_cost', 'actual_cost', 'earned_value', 'percent_complete', 'resource']),
        'risks': pd.DataFrame(columns=['risk_id', 'project_id', 'risk_description', 'probability', 'impact_cost', 'risk_score', 'mitigation_plan', 'risk_status', 'owner', 'collaboration_link']),
        'changes': pd.DataFrame(columns=['change_id', 'project_id', 'description', 'impact_cost', 'impact_schedule_days', 'status', 'date_submitted', 'collaboration_link']),
        'history': pd.DataFrame(columns=['project_id', 'month', 'cpi', 'spi']),
        'benefits': pd.DataFrame(columns=['Month', 'ReportingTimeSaved_hrs', 'CostOverrunsAvoided_k', 'ForecastAccuracy_perc'])
    }

    if use_mock_data:
        try:
            dfs['projects'], dfs['tasks'], dfs['risks'], dfs['changes'], dfs['history'], dfs['benefits'] = generate_mock_data()
            for key in empty_dfs: # Ensure all keys exist
                if key not in dfs: dfs[key] = empty_dfs[key].copy()
            data_loaded_successfully = True
        except Exception as e:
            logging.error(f"Fatal error generating mock data: {e}", exc_info=True)
            st.error(f"Fatal error generating mock data: {e}")
            dfs = {k: v.copy() for k, v in empty_dfs.items()}
            return dfs, False, ["Mock data generation failed."]
    else:
        # --- Production Data Loading Logic Placeholder ---
        logging.info("Loading data from production sources (Placeholder)...")
        st.info("Connecting to real data sources... (This is a placeholder)")
        # Simulate loading - replace with actual DB/API calls
        time.sleep(1)
        try:
            # Example: Load Projects (replace with real logic)
            # db_conn = connect_to_database()
            # dfs['projects'] = pd.read_sql("SELECT ... FROM projects", db_conn)
            dfs['projects'] = generate_mock_data(num_projects=20)[0] # Placeholder
            logging.info("Simulated loading projects.")
            # Example: Load Tasks (replace with real logic)
            # tasks_json = fetch_from_api("...")
            # dfs['tasks'] = pd.json_normalize(tasks_json['data'])
            dfs['tasks'] = generate_mock_data(num_projects=20)[1] # Placeholder
            logging.info("Simulated loading tasks.")
            # Load others similarly (using mock as placeholder for now)
            _, _, dfs['risks'], dfs['changes'], dfs['history'], dfs['benefits'] = generate_mock_data(num_projects=20)
            for key in empty_dfs: # Ensure all keys exist
                if key not in dfs: dfs[key] = empty_dfs[key].copy()
            data_loaded_successfully = True # Assume success for placeholder
        except Exception as e:
             logging.error(f"Error loading production data (Placeholder): {e}", exc_info=True)
             validation_errors.append(f"Failed to load production data: {e}")
             dfs = {k: v.copy() for k, v in empty_dfs.items()} # Return empty on error
             return dfs, False, validation_errors
        # --- End Production Logic ---

    # --- Data Validation Step ---
    if data_loaded_successfully:
        final_validation_errors = []
        validation_passed_all = True
        for name, df in dfs.items():
            is_valid, msg = validate_data(df, name)
            if not is_valid:
                final_validation_errors.append(msg)
                if name in ['projects', 'tasks']: # Essential data
                    validation_passed_all = False
                    logging.error(f"Fatal validation error for {name}: {msg}.")
                else:
                    logging.warning(f"Non-fatal validation issue for {name}: {msg}")

        if final_validation_errors:
            log_msg = f"Data validation issues: {'; '.join(final_validation_errors)}"
            if not validation_passed_all:
                logging.error(log_msg + " (Includes fatal errors)")
                st.error("Critical data validation failed. App may not function correctly.")
                return {k: v.copy() for k, v in empty_dfs.items()}, False, final_validation_errors
            else:
                logging.warning(log_msg + " (Non-fatal issues)")
                st.warning(f"Data validation issues found: {'; '.join(final_validation_errors)}. Some features might be affected.")

        if validation_passed_all:
             logging.info("All critical data passed validation.")
             data_loaded_successfully = True
        else:
             data_loaded_successfully = False # Mark as failed if critical validation failed

    # --- Data Cleaning / Transformation ---
    if data_loaded_successfully:
        try:
            logging.info("Performing data type conversions and cleaning...")
            # Apply type conversions safely, checking column existence first
            for df_name, df_content in dfs.items():
                if df_name == 'projects':
                    if 'budget' in df_content.columns: df_content['budget'] = pd.to_numeric(df_content['budget'], errors='coerce').fillna(0.0)
                    if 'start_date' in df_content.columns: df_content['start_date'] = pd.to_datetime(df_content['start_date'], errors='coerce')
                    if 'planned_end_date' in df_content.columns: df_content['planned_end_date'] = pd.to_datetime(df_content['planned_end_date'], errors='coerce')
                    if 'target_cpi' in df_content.columns: df_content['target_cpi'] = pd.to_numeric(df_content['target_cpi'], errors='coerce').fillna(1.0)
                    if 'target_spi' in df_content.columns: df_content['target_spi'] = pd.to_numeric(df_content['target_spi'], errors='coerce').fillna(1.0)
                elif df_name == 'tasks':
                    if 'planned_cost' in df_content.columns: df_content['planned_cost'] = pd.to_numeric(df_content['planned_cost'], errors='coerce').fillna(0.0)
                    if 'actual_cost' in df_content.columns: df_content['actual_cost'] = pd.to_numeric(df_content['actual_cost'], errors='coerce').fillna(0.0)
                    if 'earned_value' in df_content.columns: df_content['earned_value'] = pd.to_numeric(df_content['earned_value'], errors='coerce').fillna(0.0)
                    if 'percent_complete' in df_content.columns: df_content['percent_complete'] = pd.to_numeric(df_content['percent_complete'], errors='coerce').fillna(0.0)
                    for col in ['planned_start', 'planned_end', 'actual_start', 'actual_end']:
                         if col in df_content.columns: df_content[col] = pd.to_datetime(df_content[col], errors='coerce')
                elif df_name == 'risks':
                    if 'probability' in df_content.columns: df_content['probability'] = pd.to_numeric(df_content['probability'], errors='coerce').fillna(0.0)
                    if 'impact_cost' in df_content.columns: df_content['impact_cost'] = pd.to_numeric(df_content['impact_cost'], errors='coerce').fillna(0.0)
                    if 'risk_score' in df_content.columns: df_content['risk_score'] = pd.to_numeric(df_content['risk_score'], errors='coerce').fillna(0.0)
                elif df_name == 'changes':
                    if 'impact_cost' in df_content.columns: df_content['impact_cost'] = pd.to_numeric(df_content['impact_cost'], errors='coerce').fillna(0.0)
                    if 'impact_schedule_days' in df_content.columns: df_content['impact_schedule_days'] = pd.to_numeric(df_content['impact_schedule_days'], errors='coerce').fillna(0.0)
                    if 'date_submitted' in df_content.columns: df_content['date_submitted'] = pd.to_datetime(df_content['date_submitted'], errors='coerce')
                elif df_name == 'history':
                    if 'cpi' in df_content.columns: df_content['cpi'] = pd.to_numeric(df_content['cpi'], errors='coerce').fillna(0.0)
                    if 'spi' in df_content.columns: df_content['spi'] = pd.to_numeric(df_content['spi'], errors='coerce').fillna(0.0)
                elif df_name == 'benefits':
                    if 'ReportingTimeSaved_hrs' in df_content.columns: df_content['ReportingTimeSaved_hrs'] = pd.to_numeric(df_content['ReportingTimeSaved_hrs'], errors='coerce').fillna(0.0)
                    if 'CostOverrunsAvoided_k' in df_content.columns: df_content['CostOverrunsAvoided_k'] = pd.to_numeric(df_content['CostOverrunsAvoided_k'], errors='coerce').fillna(0.0)
                    if 'ForecastAccuracy_perc' in df_content.columns: df_content['ForecastAccuracy_perc'] = pd.to_numeric(df_content['ForecastAccuracy_perc'], errors='coerce').fillna(0.0)
            logging.info("Data type conversions and cleaning complete.")
        except Exception as e:
            logging.error(f"Error during data cleaning/conversion: {e}", exc_info=True)
            st.error("An error occurred during data preparation. Some data might be incorrect.")
            data_loaded_successfully = False
            return {k: v.copy() for k, v in empty_dfs.items()}, False, ["Data cleaning failed."]

    return dfs, data_loaded_successfully, validation_errors


# --- KPI Calculation Functions ---
def calculate_project_kpis(tasks_data, project_budget):
    """Calculates key performance indicators for a single project."""
    kpi_results = {'spi': 0.0, 'cpi': 0.0, 'pv': 0.0, 'ev': 0.0, 'ac': 0.0, 'cv': 0.0, 'sv': 0.0, 'bac': 0.0, 'eac_cpi': 0.0, 'etc_cpi': 0.0, 'vac': 0.0}
    kpi_results['bac'] = pd.to_numeric(project_budget, errors='coerce')
    if pd.isna(kpi_results['bac']): kpi_results['bac'] = 0.0
    if tasks_data is None or tasks_data.empty: return kpi_results
    required_cols = ['planned_cost', 'actual_cost', 'earned_value', 'planned_start']
    if not all(col in tasks_data.columns for col in required_cols):
        missing = [col for col in required_cols if col not in tasks_data.columns]
        logging.error(f"KPI Calc: Missing required task columns: {missing}")
        return kpi_results
    try:
        tasks_data = tasks_data.copy()
        tasks_data['planned_cost'] = pd.to_numeric(tasks_data['planned_cost'], errors='coerce').fillna(0.0)
        tasks_data['actual_cost'] = pd.to_numeric(tasks_data['actual_cost'], errors='coerce').fillna(0.0)
        tasks_data['earned_value'] = pd.to_numeric(tasks_data['earned_value'], errors='coerce').fillna(0.0)
        tasks_data['planned_start'] = pd.to_datetime(tasks_data['planned_start'], errors='coerce')
        today = pd.Timestamp(datetime.date.today())
        valid_start_dates = tasks_data.dropna(subset=['planned_start'])
        pv = valid_start_dates.loc[valid_start_dates['planned_start'] <= today, 'planned_cost'].sum()
        ev = tasks_data['earned_value'].sum()
        ac = tasks_data['actual_cost'].sum()
        bac = kpi_results['bac']
        cpi = (ev / ac) if ac != 0 else 0.0
        spi = (ev / pv) if pv != 0 else 0.0
        cv = ev - ac; sv = ev - pv
        eac_cpi = (bac / cpi) if cpi > 0.01 else bac # Use BAC if CPI is near zero to avoid extreme EAC
        etc_cpi = (eac_cpi - ac)
        vac = (bac - eac_cpi)
        kpi_results.update({'spi': spi, 'cpi': cpi, 'pv': pv, 'ev': ev, 'ac': ac, 'cv': cv, 'sv': sv, 'bac': bac, 'eac_cpi': eac_cpi, 'etc_cpi': etc_cpi, 'vac': vac})
    except Exception as e:
        logging.error(f"Error calculating KPIs for budget {project_budget}: {e}", exc_info=True)
        return {'spi': 0.0, 'cpi': 0.0, 'pv': 0.0, 'ev': 0.0, 'ac': 0.0, 'cv': 0.0, 'sv': 0.0, 'bac': kpi_results['bac'], 'eac_cpi': 0.0, 'etc_cpi': 0.0, 'vac': 0.0}
    return kpi_results

def calculate_portfolio_kpis(filtered_projects_df, filtered_tasks_df):
    """Calculates KPIs for the entire filtered portfolio."""
    logging.info("Calculating portfolio KPIs...")
    portfolio_budget = 0.0
    if 'budget' in filtered_projects_df.columns and not filtered_projects_df.empty:
         portfolio_budget = pd.to_numeric(filtered_projects_df['budget'], errors='coerce').fillna(0.0).sum()
    portfolio_kpis = calculate_project_kpis(filtered_tasks_df, portfolio_budget)
    logging.info("Portfolio KPIs calculation complete.")
    return portfolio_kpis

def add_kpis_to_projects(projects_df, tasks_df):
    """Calculates KPIs for each project and merges them into the projects dataframe."""
    logging.info("Calculating KPIs for individual projects...")
    kpi_cols_to_add = ['spi', 'cpi', 'pv', 'ev', 'ac', 'cv', 'sv', 'bac', 'eac_cpi', 'etc_cpi', 'vac']
    projects_df_out = projects_df.copy()
    for col in kpi_cols_to_add:
        if col not in projects_df_out.columns: projects_df_out[col] = 0.0
    if projects_df_out.empty or 'project_id' not in projects_df_out.columns:
        logging.warning("Cannot calculate project KPIs: Projects empty or missing 'project_id'.")
        return projects_df_out
    if tasks_df.empty or 'project_id' not in tasks_df.columns:
        logging.warning("Task data empty or missing 'project_id'. Project KPIs (except BAC) will be zero.")
        if 'budget' in projects_df_out.columns: projects_df_out['bac'] = pd.to_numeric(projects_df_out['budget'], errors='coerce').fillna(0.0)
        else: projects_df_out['bac'] = 0.0
        return projects_df_out
    project_kpis_list = []
    if 'budget' not in projects_df_out.columns:
        logging.warning("Missing 'budget' column. BAC will be zero.")
        projects_df_out['budget'] = 0.0
    else: projects_df_out['budget'] = pd.to_numeric(projects_df_out['budget'], errors='coerce').fillna(0.0)
    for pid in projects_df_out['project_id'].unique():
        proj_tasks = tasks_df[tasks_df['project_id'] == pid]
        proj_budget_series = projects_df_out.loc[projects_df_out['project_id'] == pid, 'budget']
        proj_budget = proj_budget_series.iloc[0] if not proj_budget_series.empty else 0.0
        kpis = calculate_project_kpis(proj_tasks, proj_budget)
        kpis['project_id'] = pid
        project_kpis_list.append(kpis)
    if not project_kpis_list:
        logging.warning("No KPIs calculated for any project.")
        if 'budget' in projects_df_out.columns: projects_df_out['bac'] = pd.to_numeric(projects_df_out['budget'], errors='coerce').fillna(0.0)
        else: projects_df_out['bac'] = 0.0
        return projects_df_out
    kpi_df = pd.DataFrame(project_kpis_list)
    merge_cols = kpi_cols_to_add + ['project_id']
    for col in merge_cols:
         if col not in kpi_df.columns: kpi_df[col] = 0.0 if col != 'project_id' else np.nan
    kpi_df = kpi_df[merge_cols]
    cols_to_drop = [col for col in kpi_cols_to_add if col in projects_df_out.columns]
    projects_df_base = projects_df_out.drop(columns=cols_to_drop)
    projects_with_kpis = projects_df_base.merge(kpi_df, on='project_id', how='left')
    for col in kpi_cols_to_add:
        if col in projects_with_kpis.columns: projects_with_kpis[col] = projects_with_kpis[col].fillna(0.0)
        else: projects_with_kpis[col] = 0.0
    logging.info("Individual project KPIs calculation and merge complete.")
    return projects_with_kpis

# --- Formatting Helpers ---
def format_currency(value, compact=False):
    """Formats a number as currency."""
    if pd.isna(value) or not isinstance(value, (int, float, np.number)): return "$0"
    try:
        value = float(value); abs_value = abs(value); sign = "-" if value < 0 else ""
        if compact:
            if abs_value >= 1_000_000_000: return f"{sign}${abs_value / 1_000_000_000:.1f}B"
            if abs_value >= 1_000_000: return f"{sign}${abs_value / 1_000_000:.1f}M"
            if abs_value >= 1_000: return f"{sign}${abs_value / 1_000:.0f}K"
            return f"{sign}${abs_value:.0f}"
        else: return f"{sign}${value:,.0f}"
    except (ValueError, TypeError): return "$0"

def create_gauge(value, title, max_value=2, target=1.0, color=ARCADIS_SECONDARY_PALETTE[0]):
    """Creates a Plotly gauge chart."""
    gauge_value = pd.to_numeric(value, errors='coerce')
    if pd.isna(gauge_value): gauge_value = 0.0
    fig = go.Figure(go.Indicator(
        mode = "gauge+number", value = gauge_value,
        title = {'text': title, 'font': {'size': 16, 'color': '#555555'}},
        number = {'valueformat': '.2f', 'font': {'size': 28, 'color': ARCADIS_BLACK}},
        gauge = { 'axis': {'range': [0, max_value], 'tickwidth': 1, 'tickcolor': "darkgrey"},
                  'bar': {'color': color, 'thickness': 0.4}, 'bgcolor': "#F8F9FA", 'borderwidth': 0,
                  'steps': [ {'range': [0, target * 0.9], 'color': '#FADBD8'}, {'range': [target * 0.9, target * 1.1], 'color': '#FEF9E7'}, {'range': [target * 1.1, max_value], 'color': '#D5F5E3'}],
                  'threshold': { 'line': {'color': ARCADIS_BLACK, 'width': 4}, 'thickness': 0.8, 'value': target}
                }))
    fig.update_layout(height=180, margin=dict(l=25, r=25, t=45, b=15), font=dict(family="Arial, sans-serif"))
    return fig


# --- Initialize Session State ---
# Use a more specific prefix to avoid collisions if multiple apps run in same session
SESSION_PREFIX = "pmo_pulse_"
default_state_pmo = {
    f'{SESSION_PREFIX}data_loaded': False, f'{SESSION_PREFIX}all_dfs': {},
    f'{SESSION_PREFIX}projects_df': pd.DataFrame(), f'{SESSION_PREFIX}tasks_df': pd.DataFrame(),
    f'{SESSION_PREFIX}risks_df': pd.DataFrame(), f'{SESSION_PREFIX}changes_df': pd.DataFrame(),
    f'{SESSION_PREFIX}history_df': pd.DataFrame(), f'{SESSION_PREFIX}benefits_df': pd.DataFrame(),
    f'{SESSION_PREFIX}filtered_projects_with_kpis': pd.DataFrame(),
    f'{SESSION_PREFIX}filtered_tasks': pd.DataFrame(), f'{SESSION_PREFIX}filtered_risks': pd.DataFrame(),
    f'{SESSION_PREFIX}filtered_changes': pd.DataFrame(), f'{SESSION_PREFIX}filtered_history': pd.DataFrame(),
    f'{SESSION_PREFIX}portfolio_kpis': {}, f'{SESSION_PREFIX}load_errors': [],
    # Filter states
    f'{SESSION_PREFIX}selected_sector': 'All', f'{SESSION_PREFIX}selected_status': 'All',
    f'{SESSION_PREFIX}selected_manager': 'All', f'{SESSION_PREFIX}selected_alignment': 'All',
    f'{SESSION_PREFIX}search_term': '',
    # Config states
    f'{SESSION_PREFIX}spi_threshold': 0.90, f'{SESSION_PREFIX}cpi_threshold': 0.90,
}
for key, value in default_state_pmo.items():
    if key not in st.session_state: st.session_state[key] = value

# --- Initial Data Load Trigger ---
if not st.session_state[f'{SESSION_PREFIX}data_loaded']:
    logging.info("Initial data load triggered.")
    use_mock_init = True # Default to mock data on first load
    all_dfs_init, data_loaded_ok, load_errors_init = load_data(use_mock_data=use_mock_init)
    st.session_state[f'{SESSION_PREFIX}all_dfs'] = all_dfs_init
    st.session_state[f'{SESSION_PREFIX}projects_df'] = all_dfs_init.get('projects', pd.DataFrame())
    st.session_state[f'{SESSION_PREFIX}tasks_df'] = all_dfs_init.get('tasks', pd.DataFrame())
    st.session_state[f'{SESSION_PREFIX}risks_df'] = all_dfs_init.get('risks', pd.DataFrame())
    st.session_state[f'{SESSION_PREFIX}changes_df'] = all_dfs_init.get('changes', pd.DataFrame())
    st.session_state[f'{SESSION_PREFIX}history_df'] = all_dfs_init.get('history', pd.DataFrame())
    st.session_state[f'{SESSION_PREFIX}benefits_df'] = all_dfs_init.get('benefits', pd.DataFrame())
    st.session_state[f'{SESSION_PREFIX}data_loaded'] = data_loaded_ok
    st.session_state[f'{SESSION_PREFIX}load_errors'] = load_errors_init
    # Calculate initial KPIs if data loaded successfully
    if data_loaded_ok:
         with st.spinner("Calculating initial KPIs..."):
            st.session_state[f'{SESSION_PREFIX}filtered_projects_with_kpis'] = add_kpis_to_projects(
                st.session_state[f'{SESSION_PREFIX}projects_df'],
                st.session_state[f'{SESSION_PREFIX}tasks_df']
            )
            st.session_state[f'{SESSION_PREFIX}portfolio_kpis'] = calculate_portfolio_kpis(
                st.session_state[f'{SESSION_PREFIX}filtered_projects_with_kpis'],
                st.session_state[f'{SESSION_PREFIX}tasks_df'] # Use all tasks for initial portfolio KPI
            )
            # Initialize filtered dataframes to the full set initially
            st.session_state[f'{SESSION_PREFIX}filtered_tasks'] = st.session_state[f'{SESSION_PREFIX}tasks_df'].copy()
            st.session_state[f'{SESSION_PREFIX}filtered_risks'] = st.session_state[f'{SESSION_PREFIX}risks_df'].copy()
            st.session_state[f'{SESSION_PREFIX}filtered_changes'] = st.session_state[f'{SESSION_PREFIX}changes_df'].copy()
            st.session_state[f'{SESSION_PREFIX}filtered_history'] = st.session_state[f'{SESSION_PREFIX}history_df'].copy()

    st.rerun() # Rerun after initial load/KPI calc


# --- Sidebar ---
st.sidebar.title("⚙️ Configuration & Filters")
st.sidebar.markdown("---")

# Data Status & Refresh
st.sidebar.subheader("1. Data Status")
if st.session_state[f'{SESSION_PREFIX}data_loaded']:
    st.sidebar.success(f"✅ Projects: {len(st.session_state[f'{SESSION_PREFIX}projects_df'])} | Tasks: {len(st.session_state[f'{SESSION_PREFIX}tasks_df'])}")
    if st.session_state[f'{SESSION_PREFIX}load_errors']:
        with st.sidebar.expander("Show Load/Validation Warnings"):
             for err in st.session_state[f'{SESSION_PREFIX}load_errors']: st.warning(err)
    if st.sidebar.button("🔄 Refresh Data (Simulated)", key="refresh_data_sb"):
        with st.spinner("Reloading data..."):
            st.cache_data.clear()
            st.success("Data cache cleared. Rerunning...")
            time.sleep(1)
            st.rerun()
else:
    st.sidebar.error("❌ Data not loaded.")
    if st.session_state[f'{SESSION_PREFIX}load_errors']:
        with st.sidebar.expander("Show Load Errors"):
             for err in st.session_state[f'{SESSION_PREFIX}load_errors']: st.error(err)
    if st.sidebar.button("🔄 Load Mock Data", key="load_mock_sb"):
        # Reset state and trigger mock data load
        st.session_state[f'{SESSION_PREFIX}data_loaded'] = False
        st.session_state[f'{SESSION_PREFIX}projects_df'] = pd.DataFrame() # Clear existing
        st.rerun()

st.sidebar.markdown("---")

# Filters
st.sidebar.subheader("2. Portfolio Filters")
# Use session state keys with prefix
st.session_state[f'{SESSION_PREFIX}search_term'] = st.sidebar.text_input("Search Project Name", value=st.session_state[f'{SESSION_PREFIX}search_term'], key="search_sb")

# Get filter options based on the *currently loaded* projects_df
projects_df_sb = st.session_state[f'{SESSION_PREFIX}projects_df']
sectors_sb = ['All'] + sorted(projects_df_sb['sector'].dropna().unique().tolist()) if 'sector' in projects_df_sb else ['All']
statuses_sb = ['All'] + sorted(projects_df_sb['status'].dropna().unique().tolist()) if 'status' in projects_df_sb else ['All']
managers_sb = ['All'] + sorted(projects_df_sb['project_manager'].dropna().unique().tolist()) if 'project_manager' in projects_df_sb else ['All']
alignments_sb = ['All'] + sorted(projects_df_sb['strategic_alignment'].dropna().unique().tolist()) if 'strategic_alignment' in projects_df_sb else ['All']

st.session_state[f'{SESSION_PREFIX}selected_sector'] = st.sidebar.selectbox("Sector", sectors_sb, index=sectors_sb.index(st.session_state[f'{SESSION_PREFIX}selected_sector']) if st.session_state[f'{SESSION_PREFIX}selected_sector'] in sectors_sb else 0, key="sector_sb")
st.session_state[f'{SESSION_PREFIX}selected_status'] = st.sidebar.selectbox("Status", statuses_sb, index=statuses_sb.index(st.session_state[f'{SESSION_PREFIX}selected_status']) if st.session_state[f'{SESSION_PREFIX}selected_status'] in statuses_sb else 0, key="status_sb")
st.session_state[f'{SESSION_PREFIX}selected_manager'] = st.sidebar.selectbox("Project Manager", managers_sb, index=managers_sb.index(st.session_state[f'{SESSION_PREFIX}selected_manager']) if st.session_state[f'{SESSION_PREFIX}selected_manager'] in managers_sb else 0, key="manager_sb")
st.session_state[f'{SESSION_PREFIX}selected_alignment'] = st.sidebar.selectbox("Strategic Alignment", alignments_sb, index=alignments_sb.index(st.session_state[f'{SESSION_PREFIX}selected_alignment']) if st.session_state[f'{SESSION_PREFIX}selected_alignment'] in alignments_sb else 0, key="alignment_sb")

st.sidebar.markdown("---")

# KPI Thresholds
st.sidebar.subheader("3. KPI Thresholds")
st.session_state[f'{SESSION_PREFIX}spi_threshold'] = st.sidebar.slider("Low SPI Factor", 0.5, 1.0, st.session_state[f'{SESSION_PREFIX}spi_threshold'], 0.01, format="%.2f", key="spi_thresh_sb", help="Flag projects with SPI < Target * Factor")
st.session_state[f'{SESSION_PREFIX}cpi_threshold'] = st.sidebar.slider("Low CPI Factor", 0.5, 1.0, st.session_state[f'{SESSION_PREFIX}cpi_threshold'], 0.01, format="%.2f", key="cpi_thresh_sb", help="Flag projects with CPI < Target * Factor")

# --- Apply Filters (Logic moved here, before tabs are rendered) ---
# Start with the full loaded data
projects_to_filter = st.session_state[f'{SESSION_PREFIX}projects_df'].copy()
tasks_to_filter = st.session_state[f'{SESSION_PREFIX}tasks_df'].copy()
risks_to_filter = st.session_state[f'{SESSION_PREFIX}risks_df'].copy()
changes_to_filter = st.session_state[f'{SESSION_PREFIX}changes_df'].copy()
history_to_filter = st.session_state[f'{SESSION_PREFIX}history_df'].copy()

# Apply search first
search_term_sb = st.session_state[f'{SESSION_PREFIX}search_term'].lower()
if search_term_sb and 'project_name' in projects_to_filter.columns:
    projects_to_filter = projects_to_filter[projects_to_filter['project_name'].astype(str).str.lower().str.contains(search_term_sb)]

# Apply dropdown filters
if st.session_state[f'{SESSION_PREFIX}selected_sector'] != 'All' and 'sector' in projects_to_filter:
    projects_to_filter = projects_to_filter[projects_to_filter['sector'] == st.session_state[f'{SESSION_PREFIX}selected_sector']]
if st.session_state[f'{SESSION_PREFIX}selected_status'] != 'All' and 'status' in projects_to_filter:
    projects_to_filter = projects_to_filter[projects_to_filter['status'] == st.session_state[f'{SESSION_PREFIX}selected_status']]
if st.session_state[f'{SESSION_PREFIX}selected_manager'] != 'All' and 'project_manager' in projects_to_filter:
    projects_to_filter = projects_to_filter[projects_to_filter['project_manager'] == st.session_state[f'{SESSION_PREFIX}selected_manager']]
if st.session_state[f'{SESSION_PREFIX}selected_alignment'] != 'All' and 'strategic_alignment' in projects_to_filter:
    projects_to_filter = projects_to_filter[projects_to_filter['strategic_alignment'] == st.session_state[f'{SESSION_PREFIX}selected_alignment']]

# Get IDs of filtered projects
filtered_project_ids = projects_to_filter['project_id'].tolist() if 'project_id' in projects_to_filter else []

# Filter related dataframes
if filtered_project_ids:
    st.session_state[f'{SESSION_PREFIX}filtered_tasks'] = tasks_to_filter[tasks_to_filter['project_id'].isin(filtered_project_ids)] if 'project_id' in tasks_to_filter else pd.DataFrame()
    st.session_state[f'{SESSION_PREFIX}filtered_risks'] = risks_to_filter[risks_to_filter['project_id'].isin(filtered_project_ids)] if 'project_id' in risks_to_filter else pd.DataFrame()
    st.session_state[f'{SESSION_PREFIX}filtered_changes'] = changes_to_filter[changes_to_filter['project_id'].isin(filtered_project_ids)] if 'project_id' in changes_to_filter else pd.DataFrame()
    st.session_state[f'{SESSION_PREFIX}filtered_history'] = history_to_filter[history_to_filter['project_id'].isin(filtered_project_ids)] if 'project_id' in history_to_filter else pd.DataFrame()
else:
    # Ensure empty if no projects match
    st.session_state[f'{SESSION_PREFIX}filtered_tasks'] = pd.DataFrame()
    st.session_state[f'{SESSION_PREFIX}filtered_risks'] = pd.DataFrame()
    st.session_state[f'{SESSION_PREFIX}filtered_changes'] = pd.DataFrame()
    st.session_state[f'{SESSION_PREFIX}filtered_history'] = pd.DataFrame()

# Calculate KPIs on the filtered projects
if not projects_to_filter.empty:
     with st.spinner("Calculating KPIs for filtered projects..."):
        st.session_state[f'{SESSION_PREFIX}filtered_projects_with_kpis'] = add_kpis_to_projects(
            projects_to_filter, # Pass the already filtered projects
            st.session_state[f'{SESSION_PREFIX}filtered_tasks'] # Pass the already filtered tasks
        )
        st.session_state[f'{SESSION_PREFIX}portfolio_kpis'] = calculate_portfolio_kpis(
            st.session_state[f'{SESSION_PREFIX}filtered_projects_with_kpis'],
            st.session_state[f'{SESSION_PREFIX}filtered_tasks']
        )
else:
     st.session_state[f'{SESSION_PREFIX}filtered_projects_with_kpis'] = pd.DataFrame()
     st.session_state[f'{SESSION_PREFIX}portfolio_kpis'] = calculate_portfolio_kpis(pd.DataFrame(), pd.DataFrame())


# --- Main Application Area ---
st.title("Arcadis PMO Pulse")
st.markdown(f"_Centralized dashboard for portfolio monitoring, analysis, and insights._")
if st.session_state[f'{SESSION_PREFIX}data_loaded']:
    st.success(f"Displaying data for **{len(st.session_state[f'{SESSION_PREFIX}filtered_projects_with_kpis'])}** projects based on current filters.")
else:
    st.warning("Data not loaded. Please check sidebar status or load mock data.")
st.markdown("---")


# --- Define Tabs (Narrative Structure) ---
# Renamed and reordered for narrative flow
tab_titles_pmo = [
    "👋 Welcome",
    "📊 Executive Summary",
    "🔍 Portfolio Analysis", # Was Analysis
    "🎯 Project Deep Dive", # Was Projects
    "💰 Benefits & ROI", # Was Benefits
    # "🚨 Alerts", # Removed Alerts as separate tab, integrated into Summary/Deep Dive
    "📝 Reports",
    "💾 Data & Settings" # Combined Data/Settings
]
tabs_pmo = st.tabs(tab_titles_pmo)
# Store tab map in session state for Part 2 access
st.session_state[f'{SESSION_PREFIX}tab_map_pmo'] = {name: tab for name, tab in zip(tab_titles_pmo, tabs_pmo)}


# --- Welcome Tab (Narrative Focus) ---
with st.session_state[f'{SESSION_PREFIX}tab_map_pmo']["👋 Welcome"]:
    st.markdown("<div class='welcome-section'>", unsafe_allow_html=True)
    st.header("The Challenge: Navigating Portfolio Complexity")
    st.markdown("""
    Managing a diverse project portfolio presents significant challenges: scattered data across systems, inconsistent performance metrics, difficulty identifying emerging risks, and time-consuming manual reporting. This often leads to reactive decision-making, missed opportunities for intervention, and ultimately, impacts project success and profitability.
    """)
    st.markdown("</div>", unsafe_allow_html=True)

    st.markdown("<div class='welcome-section'>", unsafe_allow_html=True)
    st.header("The Solution: PMO Pulse")
    st.markdown(f"""
    **PMO Pulse**, powered by Arcadis, provides a unified, data-driven view of your project portfolio. By integrating key data sources and automating performance calculations (like SPI and CPI), it delivers actionable insights to proactively manage projects, mitigate risks, and optimize resource allocation.
    """)
    st.markdown("</div>", unsafe_allow_html=True)

    st.markdown("<div class='welcome-section'>", unsafe_allow_html=True)
    st.header("Key Capabilities")
    col1_cap_pmo, col2_cap_pmo = st.columns(2)
    with col1_cap_pmo:
        st.markdown("""
        <ul class="styled-list">
            <li><b>Centralized Dashboard:</b> At-a-glance view of overall portfolio health and key metrics (<i>Executive Summary</i>).</li>
            <li><b>Performance Monitoring:</b> Track Cost Performance Index (CPI) and Schedule Performance Index (SPI) trends.</li>
            <li><b>Portfolio Segmentation:</b> Filter and analyze performance by sector, status, PM, etc. (<i>Portfolio Analysis</i>).</li>
            <li><b>Risk Exposure Analysis:</b> Identify high-risk projects based on aggregated risk scores (<i>Portfolio Analysis</i>).</li>
        </ul>
        """, unsafe_allow_html=True)
    with col2_cap_pmo:
        st.markdown("""
        <ul class="styled-list">
            <li><b>Detailed Project View:</b> Drill down into individual project KPIs, tasks, risks, and changes (<i>Project Deep Dive</i>).</li>
            <li><b>Change Control Monitoring:</b> Analyze the impact of change requests on cost and schedule (<i>Portfolio Analysis</i>).</li>
            <li><b>Benefits Tracking:</b> Visualize the value delivered through improved processes (<i>Benefits & ROI</i>).</li>
            <li><b>Customizable Reporting:</b> Export data and generate summary reports (<i>Reports</i>, <i>Data & Settings</i>).</li>
        </ul>
        """, unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

    st.markdown("<div class='welcome-section'>", unsafe_allow_html=True)
    st.header("Answering Your PMO Questions")
    st.markdown("""
    PMO Pulse helps answer critical questions:
    <ul class="styled-list">
        <li>❓ <b>What is the overall health of my portfolio?</b> (See <i>Executive Summary</i>)</li>
        <li>❓ <b>Which projects are falling behind schedule or over budget?</b> (Check KPIs in <i>Executive Summary</i> & <i>Portfolio Analysis</i>)</li>
        <li>❓ <b>Are there performance trends across specific sectors or managers?</b> (Filter in <i>Portfolio Analysis</i>)</li>
        <li>❓ <b>Which projects carry the highest risk exposure?</b> (Analyze in <i>Portfolio Analysis</i> & <i>Project Deep Dive</i>)</li>
        <li>❓ <b>What is the cumulative impact of change requests?</b> (See <i>Portfolio Analysis</i>)</li>
    </ul>
    <b>Use the sidebar to configure filters and explore the different analysis tabs.</b>
    """, unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

# --- Placeholder for Remaining Tabs (Implemented in Part 2 & 3) ---
# Add placeholders to avoid errors if Part 1 is run standalone
# These will be replaced by the actual content in the next parts.
placeholder_tabs = ["📊 Executive Summary", "🔍 Portfolio Analysis", "🎯 Project Deep Dive", "💰 Benefits & ROI", "📝 Reports", "💾 Data & Settings"]
for tab_name in placeholder_tabs:
     # Check if tab_map exists and the key is present before accessing
     if f'{SESSION_PREFIX}tab_map_pmo' in st.session_state and tab_name in st.session_state[f'{SESSION_PREFIX}tab_map_pmo']:
         with st.session_state[f'{SESSION_PREFIX}tab_map_pmo'][tab_name]:
             st.header(tab_name) # Use the actual tab name as header
             st.info(f"Content for {tab_name} will be implemented in the next part.")


# Note: End of Part 1. The detailed content for tabs 2-7 will be in Parts 2 & 3.
# -*- coding: utf-8 -*-
"""
PMO Pulse Narrative Driven App (v1.0 - Part 2)

Implements the Executive Summary and Portfolio Analysis tabs.
Relies on the structure and functions defined in Part 1.
"""

# ==============================================================================
# Part 2: Tab Implementations (Executive Summary, Portfolio Analysis)
# ==============================================================================

import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import datetime
from dateutil.relativedelta import relativedelta
import time
import logging
import io
import os

# --- Assume Helper Functions & Constants are available from Part 1 ---
# (e.g., format_currency, create_gauge, SESSION_PREFIX, COLOR_*, PLOTLY_TEMPLATE)
# --- Assume Session State is initialized and populated from Part 1 ---

# --- Helper Functions (If specific to these tabs, define here) ---
# (None needed specifically for Part 2 beyond what's in Part 1)


# --- Retrieve necessary variables from session state ---
# Use .get() with defaults for safety
SESSION_PREFIX = "pmo_pulse_" # Ensure prefix is defined
tab_map_pmo = st.session_state.get(f'{SESSION_PREFIX}tab_map_pmo', {})
data_loaded = st.session_state.get(f'{SESSION_PREFIX}data_loaded', False)
projects_df = st.session_state.get(f'{SESSION_PREFIX}projects_df', pd.DataFrame())
tasks_df = st.session_state.get(f'{SESSION_PREFIX}tasks_df', pd.DataFrame())
risks_df = st.session_state.get(f'{SESSION_PREFIX}risks_df', pd.DataFrame())
changes_df = st.session_state.get(f'{SESSION_PREFIX}changes_df', pd.DataFrame())
history_df = st.session_state.get(f'{SESSION_PREFIX}history_df', pd.DataFrame())
benefits_df = st.session_state.get(f'{SESSION_PREFIX}benefits_df', pd.DataFrame())
filtered_projects_with_kpis = st.session_state.get(f'{SESSION_PREFIX}filtered_projects_with_kpis', pd.DataFrame())
filtered_tasks = st.session_state.get(f'{SESSION_PREFIX}filtered_tasks', pd.DataFrame())
filtered_risks = st.session_state.get(f'{SESSION_PREFIX}filtered_risks', pd.DataFrame())
filtered_changes = st.session_state.get(f'{SESSION_PREFIX}filtered_changes', pd.DataFrame())
filtered_history = st.session_state.get(f'{SESSION_PREFIX}filtered_history', pd.DataFrame())
portfolio_kpis = st.session_state.get(f'{SESSION_PREFIX}portfolio_kpis', {})
spi_threshold = st.session_state.get(f'{SESSION_PREFIX}spi_threshold', 0.90)
cpi_threshold = st.session_state.get(f'{SESSION_PREFIX}cpi_threshold', 0.90)

# --- Plotly Template (Define if not carried over from Part 1) ---
ARCADIS_ORANGE = "#E67300"
ARCADIS_BLACK = "#000000"
ARCADIS_GREY = "#6c757d"
ARCADIS_WHITE = "#FFFFFF"
ARCADIS_DARK_GREY = "#646469"
ARCADIS_SECONDARY_PALETTE = ["#00A3A1", ARCADIS_DARK_GREY, "#D6D6D8"] # Teal, Dark Grey, Mid Grey
COLOR_SUCCESS = "#2ECC71"
COLOR_WARNING = "#F1C40F"
COLOR_DANGER = "#E74C3C"
COLOR_INFO = "#3498DB"
PLOTLY_TEMPLATE = "plotly_white" # Base template

# --- Formatting Helpers (Define if not carried over from Part 1) ---
def format_currency(value, compact=False):
    """Formats a number as currency."""
    if pd.isna(value) or not isinstance(value, (int, float, np.number)): return "$0"
    try:
        value = float(value); abs_value = abs(value); sign = "-" if value < 0 else ""
        if compact:
            if abs_value >= 1_000_000_000: return f"{sign}${abs_value / 1_000_000_000:.1f}B"
            if abs_value >= 1_000_000: return f"{sign}${abs_value / 1_000_000:.1f}M"
            if abs_value >= 1_000: return f"{sign}${abs_value / 1_000:.0f}K"
            return f"{sign}${abs_value:.0f}"
        else: return f"{sign}${value:,.0f}"
    except (ValueError, TypeError): return "$0"

def create_gauge(value, title, max_value=2, target=1.0, color=ARCADIS_SECONDARY_PALETTE[0]):
    """Creates a Plotly gauge chart."""
    gauge_value = pd.to_numeric(value, errors='coerce')
    if pd.isna(gauge_value): gauge_value = 0.0
    fig = go.Figure(go.Indicator(
        mode = "gauge+number", value = gauge_value,
        title = {'text': title, 'font': {'size': 16, 'color': '#555555'}},
        number = {'valueformat': '.2f', 'font': {'size': 28, 'color': ARCADIS_BLACK}},
        gauge = { 'axis': {'range': [0, max_value], 'tickwidth': 1, 'tickcolor': "darkgrey"},
                  'bar': {'color': color, 'thickness': 0.4}, 'bgcolor': "#F8F9FA", 'borderwidth': 0,
                  'steps': [ {'range': [0, target * 0.9], 'color': '#FADBD8'}, {'range': [target * 0.9, target * 1.1], 'color': '#FEF9E7'}, {'range': [target * 1.1, max_value], 'color': '#D5F5E3'}],
                  'threshold': { 'line': {'color': ARCADIS_BLACK, 'width': 4}, 'thickness': 0.8, 'value': target}
                }))
    fig.update_layout(height=180, margin=dict(l=25, r=25, t=45, b=15), font=dict(family="Arial, sans-serif"))
    return fig

# --- Executive Summary Tab ---
summary_tab_key = "📊 Executive Summary"
if summary_tab_key in tab_map_pmo:
    with tab_map_pmo[summary_tab_key]:
        st.header("📊 Executive Summary")
        st.markdown("_High-level overview of the portfolio's performance based on current filters._")

        if not data_loaded or filtered_projects_with_kpis.empty:
            st.info("ℹ️ No projects match the current filters or data is not loaded. Adjust filters in the sidebar or load data.")
        else:
            # --- Key Callouts / Alerts ---
            st.subheader("🚨 Portfolio Status Alerts")
            at_risk_count = len(filtered_projects_with_kpis[filtered_projects_with_kpis['status'] == 'At Risk'])
            delayed_count = len(filtered_projects_with_kpis[filtered_projects_with_kpis['status'] == 'Delayed'])
            neg_vac_count = len(filtered_projects_with_kpis[pd.to_numeric(filtered_projects_with_kpis['vac'], errors='coerce') < 0])

            alert_cols = st.columns(3)
            with alert_cols[0]:
                if delayed_count > 0: st.markdown(f"<div class='callout callout-danger'><i class='fas fa-clock'></i><div><h4>{delayed_count} Delayed Projects</h4><p>Immediate schedule review required.</p></div></div>", unsafe_allow_html=True)
                else: st.markdown(f"<div class='callout callout-success'><i class='fas fa-check-circle'></i><div><h4>0 Delayed Projects</h4><p>Schedule appears under control.</p></div></div>", unsafe_allow_html=True)
            with alert_cols[1]:
                if at_risk_count > 0: st.markdown(f"<div class='callout callout-warning'><i class='fas fa-exclamation-triangle'></i><div><h4>{at_risk_count} Projects At Risk</h4><p>Monitor closely, review mitigations.</p></div></div>", unsafe_allow_html=True)
                else: st.markdown(f"<div class='callout callout-success'><i class='fas fa-check-circle'></i><div><h4>0 Projects At Risk</h4><p>Risks appear managed.</p></div></div>", unsafe_allow_html=True)
            with alert_cols[2]:
                if neg_vac_count > 0: st.markdown(f"<div class='callout callout-danger'><i class='fas fa-dollar-sign'></i><div><h4>{neg_vac_count} Projects Forecast Over Budget</h4><p>Investigate cost drivers and EAC.</p></div></div>", unsafe_allow_html=True)
                else: st.markdown(f"<div class='callout callout-success'><i class='fas fa-check-circle'></i><div><h4>0 Projects Forecast Over Budget</h4><p>Cost forecasts look favorable.</p></div></div>", unsafe_allow_html=True)

            st.divider()

            # --- Top KPIs & Financial Forecast ---
            st.subheader("📈 Overall Performance & Forecast")
            with st.container():
                col1, col2, col3 = st.columns([2, 2, 3], gap="medium")
                with col1:
                    spi_gauge = create_gauge(portfolio_kpis.get('spi', 0), "Overall SPI", color=ARCADIS_SECONDARY_PALETTE[0]) # Teal
                    st.plotly_chart(spi_gauge, use_container_width=True)
                with col2:
                    cpi_gauge = create_gauge(portfolio_kpis.get('cpi', 0), "Overall CPI", color=ARCADIS_ORANGE)
                    st.plotly_chart(cpi_gauge, use_container_width=True)
                with col3:
                    st.markdown("##### Financial Forecast (Portfolio)")
                    fin_cols = st.columns(3)
                    with fin_cols[0]: st.metric("BAC", format_currency(portfolio_kpis.get('bac', 0), compact=True))
                    with fin_cols[1]: st.metric("EAC", format_currency(portfolio_kpis.get('eac_cpi', 0), compact=True))
                    with fin_cols[2]: st.metric("VAC", format_currency(portfolio_kpis.get('vac', 0), compact=True))

            st.divider()

            # --- KPI Trends ---
            with st.container():
                st.subheader("⏱️ Performance Trends (Portfolio Average)")
                if not filtered_history.empty and 'month' in filtered_history.columns and 'spi' in filtered_history.columns and 'cpi' in filtered_history.columns:
                    try:
                        history_numeric = filtered_history.copy()
                        history_numeric['spi'] = pd.to_numeric(history_numeric['spi'], errors='coerce')
                        history_numeric['cpi'] = pd.to_numeric(history_numeric['cpi'], errors='coerce')
                        history_numeric = history_numeric.dropna(subset=['spi', 'cpi', 'month'])

                        if not history_numeric.empty:
                            monthly_trends = history_numeric.groupby('month')[['spi', 'cpi']].mean().reset_index().sort_values('month')
                            fig_trends = go.Figure()
                            fig_trends.add_trace(go.Scatter(x=monthly_trends['month'], y=monthly_trends['spi'], mode='lines+markers', name='Average SPI', line=dict(color=ARCADIS_SECONDARY_PALETTE[0], width=2.5), marker=dict(size=7)))
                            fig_trends.add_trace(go.Scatter(x=monthly_trends['month'], y=monthly_trends['cpi'], mode='lines+markers', name='Average CPI', line=dict(color=ARCADIS_ORANGE, width=2.5), marker=dict(size=7)))
                            fig_trends.add_hline(y=1.0, line_dash="dash", line_color=ARCADIS_GREY, annotation_text="Target", annotation_position="bottom right")
                            fig_trends.update_layout(xaxis_title=None, yaxis_title="Index Value", height=350, margin=dict(t=20, b=40, l=40, r=20), legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1, font=dict(size=12)), template=PLOTLY_TEMPLATE, hovermode="x unified")
                            st.plotly_chart(fig_trends, use_container_width=True)
                        else:
                             st.info("ℹ️ No valid historical KPI data available for trends after cleaning.")
                    except Exception as e:
                        logging.error(f"Error creating KPI trend chart: {e}", exc_info=True)
                        st.warning("⚠️ Could not generate KPI trend chart.")
                else: st.info("ℹ️ No historical KPI data available for trends in the filtered data.")

            st.divider()

            # --- Projects Requiring Attention ---
            with st.container():
                st.subheader("⚠️ Projects Requiring Attention")
                # Use thresholds from session state
                spi_threshold_factor = st.session_state[f'{SESSION_PREFIX}spi_threshold']
                cpi_threshold_factor = st.session_state[f'{SESSION_PREFIX}cpi_threshold']

                attention_spi_df = pd.DataFrame()
                if all(col in filtered_projects_with_kpis.columns for col in ['spi', 'target_spi', 'status', 'project_name']):
                    numeric_spi = pd.to_numeric(filtered_projects_with_kpis['spi'], errors='coerce')
                    numeric_target_spi = pd.to_numeric(filtered_projects_with_kpis['target_spi'], errors='coerce')
                    valid_spi_comparison = filtered_projects_with_kpis[numeric_spi.notna() & numeric_target_spi.notna()]
                    attention_spi_df = valid_spi_comparison[
                        (numeric_spi > 0) &
                        (numeric_spi < numeric_target_spi * spi_threshold_factor) &
                        (valid_spi_comparison['status'] != 'Completed')
                    ].sort_values('spi')
                else: logging.warning("Missing columns for SPI attention list.")

                attention_cpi_df = pd.DataFrame()
                if all(col in filtered_projects_with_kpis.columns for col in ['cpi', 'target_cpi', 'status', 'project_name']):
                     numeric_cpi = pd.to_numeric(filtered_projects_with_kpis['cpi'], errors='coerce')
                     numeric_target_cpi = pd.to_numeric(filtered_projects_with_kpis['target_cpi'], errors='coerce')
                     valid_cpi_comparison = filtered_projects_with_kpis[numeric_cpi.notna() & numeric_target_cpi.notna()]
                     attention_cpi_df = valid_cpi_comparison[
                         (numeric_cpi > 0) &
                         (numeric_cpi < numeric_target_cpi * cpi_threshold_factor) &
                         (valid_cpi_comparison['status'] != 'Completed')
                     ].sort_values('cpi')
                else: logging.warning("Missing columns for CPI attention list.")

                attention_risk_df = pd.DataFrame()
                if not filtered_risks.empty and all(col in filtered_risks.columns for col in ['project_id', 'risk_score', 'risk_status']):
                    if 'project_id' in filtered_projects_with_kpis.columns and 'project_name' in filtered_projects_with_kpis.columns:
                        try:
                            risks_copy_att = filtered_risks.copy()
                            risks_copy_att['risk_score'] = pd.to_numeric(risks_copy_att['risk_score'], errors='coerce').fillna(0)
                            project_risk_scores = risks_copy_att[risks_copy_att['risk_status'].isin(['Open', 'Mitigating'])].groupby('project_id')['risk_score'].sum().reset_index()
                            attention_risk_df = filtered_projects_with_kpis[['project_id', 'project_name', 'status']].merge(project_risk_scores, on='project_id', how='inner').nlargest(5, 'risk_score')
                        except Exception as e: logging.error(f"Error calculating top risk projects: {e}", exc_info=True)
                    else: logging.warning("Cannot calculate top risk projects: Missing columns.")

                col1_att, col2_att, col3_att = st.columns(3, gap="medium")
                with col1_att:
                    st.markdown(f"##### <i class='fas fa-tachometer-alt'></i> Low SPI (< {spi_threshold_factor:.0%} Target)", unsafe_allow_html=True)
                    if not attention_spi_df.empty:
                        for index, row in attention_spi_df.head(5).iterrows():
                            st.markdown(f"<div class='attention-list-item spi'><strong>{row.get('project_name','N/A')}</strong> (SPI: {row.get('spi', 0):.2f} / Target: {row.get('target_spi', 0):.2f})<span>Status: {row.get('status','N/A')}</span></div>", unsafe_allow_html=True)
                    else: st.success("✅ No projects")
                with col2_att:
                    st.markdown(f"##### <i class='fas fa-coins'></i> Low CPI (< {cpi_threshold_factor:.0%} Target)", unsafe_allow_html=True)
                    if not attention_cpi_df.empty:
                        for index, row in attention_cpi_df.head(5).iterrows():
                            st.markdown(f"<div class='attention-list-item cpi'><strong>{row.get('project_name','N/A')}</strong> (CPI: {row.get('cpi', 0):.2f} / Target: {row.get('target_cpi', 0):.2f})<span>Status: {row.get('status','N/A')}</span></div>", unsafe_allow_html=True)
                    else: st.success("✅ No projects")
                with col3_att:
                    st.markdown("##### <i class='fas fa-shield-alt'></i> Top 5 by Risk Exposure", unsafe_allow_html=True)
                    if not attention_risk_df.empty:
                        for index, row in attention_risk_df.iterrows():
                            st.markdown(f"<div class='attention-list-item risk'><strong>{row.get('project_name','N/A')}</strong> (Score: {format_currency(row.get('risk_score', 0), compact=True)})<span>Status: {row.get('status','N/A')}</span></div>", unsafe_allow_html=True)
                    else: st.success("✅ No projects")

# --- Portfolio Analysis Tab ---
analysis_tab_key = "🔍 Portfolio Analysis"
if analysis_tab_key in tab_map_pmo:
    with tab_map_pmo[analysis_tab_key]:
        st.header("🔍 Portfolio Analysis")
        st.markdown("_Dive deeper into specific performance dimensions across the filtered portfolio._")

        if not data_loaded or filtered_projects_with_kpis.empty:
            st.info("ℹ️ No projects match the current filters or data is not loaded. Adjust filters in the sidebar or load data.")
        else:
            # Use the filtered dataframes already calculated and stored in session state
            projects_analysis = st.session_state[f'{SESSION_PREFIX}filtered_projects_with_kpis']
            tasks_analysis = st.session_state[f'{SESSION_PREFIX}filtered_tasks']
            risks_analysis = st.session_state[f'{SESSION_PREFIX}filtered_risks']
            changes_analysis = st.session_state[f'{SESSION_PREFIX}filtered_changes']

            analysis_options = st.selectbox("Select Analysis Area:", ["Cost & Financials", "Schedule & Progress", "Risk Exposure", "Resource Allocation", "Change Control"], key="portfolio_analysis_select")
            st.divider()

            # --- Cost & Financials ---
            if analysis_options == "Cost & Financials":
                st.markdown("### Cost Variance & Forecasting")
                required_cols_perf = ['cpi', 'spi', 'budget', 'project_name', 'sector', 'status']
                required_cols_treemap = ['sector', 'status', 'budget', 'project_name']
                required_cols_eac_vac = ['project_name', 'budget', 'eac_cpi', 'vac']

                if all(col in projects_analysis.columns for col in required_cols_perf):
                    st.markdown("#### Project Performance Matrix (CPI vs SPI)")
                    try:
                        active_projects = projects_analysis[projects_analysis['status'] != 'Completed'].copy()
                        active_projects['budget_size'] = pd.to_numeric(active_projects['budget'], errors='coerce').fillna(1.0).clip(lower=1.0)
                        if not active_projects.empty:
                            active_projects['spi'] = pd.to_numeric(active_projects['spi'], errors='coerce')
                            active_projects['cpi'] = pd.to_numeric(active_projects['cpi'], errors='coerce')
                            active_projects_plot = active_projects.dropna(subset=['spi', 'cpi'])
                            if not active_projects_plot.empty:
                                fig_perf_matrix = px.scatter(active_projects_plot, x='spi', y='cpi', size='budget_size', color='sector',
                                                            hover_name='project_name', hover_data={'spi':':.2f', 'cpi':':.2f', 'budget':':,', 'budget_size': False},
                                                            title="Active Project Performance: CPI vs. SPI (Size by Budget)", labels={'spi': 'Schedule Performance Index (SPI)', 'cpi': 'Cost Performance Index (CPI)', 'budget_size': 'Budget'},
                                                            color_discrete_sequence=px.colors.qualitative.Vivid)
                                fig_perf_matrix.add_vline(x=1.0, line_dash="dash", line_color=ARCADIS_GREY)
                                fig_perf_matrix.add_hline(y=1.0, line_dash="dash", line_color=ARCADIS_GREY)
                                fig_perf_matrix.update_layout(template=PLOTLY_TEMPLATE, height=500, legend_title_text='Sector')
                                st.plotly_chart(fig_perf_matrix, use_container_width=True)
                                st.markdown("<p class='insight-text'><i class='fas fa-lightbulb'></i><b>Insight:</b> Ideal projects are top-right (SPI > 1, CPI > 1). Bottom-left projects need urgent attention.</p>", unsafe_allow_html=True)
                            else: st.info("ℹ️ No active projects with valid SPI/CPI data found.")
                        else: st.info("ℹ️ No active (non-completed) projects match the filters.")
                    except Exception as e: logging.error(f"Error creating performance matrix chart: {e}", exc_info=True); st.warning("⚠️ Could not generate Performance Matrix chart.")
                else: st.warning(f"⚠️ Missing columns for Performance Matrix chart.")

                st.divider()
                if all(col in projects_analysis.columns for col in required_cols_treemap):
                     st.markdown("#### Budget Allocation by Sector & Status")
                     try:
                        df_treemap = projects_analysis.copy()
                        df_treemap['budget'] = pd.to_numeric(df_treemap['budget'], errors='coerce').fillna(0.0)
                        df_treemap = df_treemap[df_treemap['budget'] > 0]
                        if not df_treemap.empty:
                            df_treemap['budget_formatted'] = df_treemap['budget'].apply(lambda x: format_currency(x))
                            fig_treemap = px.treemap(df_treemap, path=[px.Constant("All Projects"), 'sector', 'status', 'project_name'], values='budget',
                                                     title='Budget Allocation Treemap', color='sector', hover_data={'budget_formatted': True, 'budget':False},
                                                     color_discrete_sequence=px.colors.qualitative.Pastel)
                            fig_treemap.update_traces(textinfo='label+percent root', hovertemplate='<b>%{label}</b><br>Budget: %{customdata[0]}<br>Percentage of Parent: %{percentParent:.1%}')
                            fig_treemap.update_layout(height=500, margin = dict(t=50, l=25, r=25, b=25))
                            st.plotly_chart(fig_treemap, use_container_width=True)
                            st.markdown("<p class='insight-text'><i class='fas fa-lightbulb'></i><b>Insight:</b> Visualize budget distribution. Large blocks representing 'At Risk' or 'Delayed' projects warrant investigation.</p>", unsafe_allow_html=True)
                        else: st.info("ℹ️ No projects with positive budget found for Treemap.")
                     except Exception as e: logging.error(f"Error creating budget treemap: {e}", exc_info=True); st.warning("⚠️ Could not generate Budget Allocation Treemap.")
                else: st.warning(f"⚠️ Missing columns for Treemap.")


                st.divider()
                if all(col in projects_analysis.columns for col in required_cols_eac_vac):
                    st.markdown("#### Budget vs. Forecast & Variance")
                    col1, col2 = st.columns(2)
                    try:
                        plot_df_eac_vac = projects_analysis.copy()
                        plot_df_eac_vac['budget'] = pd.to_numeric(plot_df_eac_vac['budget'], errors='coerce').fillna(0)
                        plot_df_eac_vac['eac_cpi'] = pd.to_numeric(plot_df_eac_vac['eac_cpi'], errors='coerce').fillna(0)
                        plot_df_eac_vac['vac'] = pd.to_numeric(plot_df_eac_vac['vac'], errors='coerce').fillna(0)
                        with col1:
                            fig_eac = go.Figure(data=[ go.Bar(name='Budget (BAC)', x=plot_df_eac_vac['project_name'], y=plot_df_eac_vac['budget'], marker_color=COLOR_PRIMARY), go.Bar(name='Forecast (EAC)', x=plot_df_eac_vac['project_name'], y=plot_df_eac_vac['eac_cpi'], marker_color=COLOR_ACCENT)])
                            fig_eac.update_layout(barmode='group', title='Budget vs. Forecast (EAC)', yaxis_title='Amount ($)', xaxis_title=None, legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1), template=PLOTLY_TEMPLATE, height=400)
                            st.plotly_chart(fig_eac, use_container_width=True)
                        with col2:
                            df_vac = plot_df_eac_vac.sort_values('vac')
                            fig_vac = px.bar(df_vac, x='project_name', y='vac', title="Variance at Completion (VAC)", labels={'project_name': 'Project', 'vac': 'VAC ($)'}, color='vac', color_continuous_scale=[COLOR_DANGER, COLOR_WARNING, COLOR_SUCCESS], text_auto=True)
                            fig_vac.update_traces(text=[format_currency(v, compact=True) for v in df_vac['vac']], textposition='outside')
                            fig_vac.update_layout(yaxis_title='Amount ($)', xaxis_title=None, template=PLOTLY_TEMPLATE, height=400, coloraxis_showscale=False)
                            st.plotly_chart(fig_vac, use_container_width=True)
                    except Exception as e: logging.error(f"Error creating EAC/VAC charts: {e}", exc_info=True); st.warning("⚠️ Could not generate Budget vs Forecast charts.")
                else: st.warning(f"⚠️ Missing columns for Budget vs Forecast charts.")

            # --- Schedule & Progress ---
            elif analysis_options == "Schedule & Progress":
                st.markdown("### Schedule Variance & Timelines")
                required_cols_sv = ['project_name', 'sv']
                required_cols_gantt_tasks = ['project_id', 'task_name', 'planned_start', 'planned_end']

                if all(col in projects_analysis.columns for col in required_cols_sv):
                    try:
                        df_sv = projects_analysis.copy()
                        df_sv['sv'] = pd.to_numeric(df_sv['sv'], errors='coerce').fillna(0)
                        df_sv = df_sv.sort_values('sv')
                        fig_sv = px.bar(df_sv, x='project_name', y='sv', title="Schedule Variance (SV = EV - PV) by Project", labels={'project_name': 'Project', 'sv': 'Schedule Variance ($ Value)'}, color='sv', color_continuous_scale=[COLOR_DANGER, COLOR_WARNING, COLOR_SUCCESS], text_auto=True)
                        fig_sv.update_traces(text=[format_currency(v, compact=True) for v in df_sv['sv']], textposition='outside')
                        fig_sv.update_layout(yaxis_title='Value ($)', xaxis_title=None, template=PLOTLY_TEMPLATE, height=400, coloraxis_showscale=False)
                        st.plotly_chart(fig_sv, use_container_width=True)
                        st.markdown("<p class='insight-text'><i class='fas fa-lightbulb'></i><b>Insight:</b> Negative SV (red/orange) indicates schedule delays in value delivery. Positive SV (green) suggests being ahead.</p>", unsafe_allow_html=True)
                    except Exception as e: logging.error(f"Error creating SV chart: {e}", exc_info=True); st.warning("⚠️ Could not generate Schedule Variance chart.")
                else: st.warning(f"⚠️ Missing columns for Schedule Variance chart.")

                st.divider()
                st.markdown("#### Project Timeline (Sample Tasks)")
                if not tasks_analysis.empty and 'project_id' in tasks_analysis.columns and 'project_id' in projects_analysis.columns and 'project_name' in projects_analysis.columns and all(col in tasks_analysis.columns for col in required_cols_gantt_tasks):
                    tasks_for_gantt_base = tasks_analysis.head(150).merge(projects_analysis[['project_id', 'project_name']], on='project_id', how='left')
                    tasks_for_gantt = tasks_for_gantt_base.dropna(subset=['project_name', 'planned_start', 'planned_end']).copy()
                    if not tasks_for_gantt.empty:
                        try:
                            tasks_for_gantt['planned_start'] = pd.to_datetime(tasks_for_gantt['planned_start'])
                            tasks_for_gantt['planned_end'] = pd.to_datetime(tasks_for_gantt['planned_end'])
                            custom_data_gantt = ["planned_start", "planned_end"]
                            if 'resource' in tasks_for_gantt.columns: custom_data_gantt.append("resource")
                            fig_gantt = px.timeline(tasks_for_gantt, x_start="planned_start", x_end="planned_end", y="project_name", color="project_name", title="Planned Timelines (Task Sample)", labels={"project_name": "Project"}, hover_name="task_name", custom_data=custom_data_gantt)
                            min_date = tasks_for_gantt['planned_start'].min() - pd.Timedelta(days=15); max_date = tasks_for_gantt['planned_end'].max() + pd.Timedelta(days=15)
                            fig_gantt.update_layout(xaxis_title="Date", yaxis_title=None, xaxis_range=[min_date, max_date], template=PLOTLY_TEMPLATE, height=500)
                            fig_gantt.update_yaxes(categoryorder='total ascending')
                            hover_template = "<b>Task:</b> %{hovertext}<br><b>Project:</b> %{y}<br>"
                            if 'resource' in tasks_for_gantt.columns: hover_template += "<b>Resource:</b> %{customdata[2]}<br>"
                            hover_template += "<b>Planned:</b> %{customdata[0]|%Y-%m-%d} to %{customdata[1]|%Y-%m-%d}<extra></extra>"
                            fig_gantt.update_traces(hovertemplate=hover_template)
                            st.plotly_chart(fig_gantt, use_container_width=True)
                            st.caption("Displaying planned timelines for a sample of up to 150 tasks.")
                        except Exception as e: logging.error(f"Error creating Gantt chart: {e}", exc_info=True); st.warning("⚠️ Could not generate Gantt chart.")
                    else: st.info("ℹ️ No valid task data with project names and dates found for Gantt chart after filtering.")
                else: st.warning("⚠️ Cannot generate Gantt chart. Task or Project data is empty or missing required columns.")


            # --- Risk Exposure ---
            elif analysis_options == "Risk Exposure":
                st.markdown("### Risk Landscape & Breakdown")
                required_cols_matrix = ['probability', 'impact_cost', 'risk_score', 'risk_description', 'project_id', 'risk_status', 'mitigation_plan', 'owner']
                required_cols_sunburst = ['owner', 'risk_status', 'risk_score']
                required_cols_table = ['project_id', 'risk_description', 'probability', 'impact_cost', 'risk_score', 'risk_status', 'mitigation_plan', 'owner']

                if not risks_analysis.empty and all(col in risks_analysis.columns for col in required_cols_matrix):
                    st.markdown("#### Risk Matrix (Probability vs. Impact Cost)")
                    try:
                        risk_matrix_df = risks_analysis.copy()
                        risk_matrix_df['probability'] = pd.to_numeric(risk_matrix_df['probability'], errors='coerce').fillna(0)
                        risk_matrix_df['impact_cost'] = pd.to_numeric(risk_matrix_df['impact_cost'], errors='coerce').fillna(0)
                        risk_matrix_df['risk_score'] = pd.to_numeric(risk_matrix_df['risk_score'], errors='coerce').fillna(1).clip(lower=1)
                        fig_risk_matrix = px.scatter(risk_matrix_df, x='probability', y='impact_cost', size='risk_score', color='risk_score', hover_name='risk_description', hover_data=['project_id', 'risk_status', 'mitigation_plan', 'owner'], title="Risk Landscape (Size & Color by Risk Score)", labels={'probability': 'Probability', 'impact_cost': 'Impact Cost ($)', 'risk_score': 'Risk Score'}, color_continuous_scale="OrRd")
                        min_y = 0; max_y = risk_matrix_df['impact_cost'].max() * 1.1 if not risk_matrix_df['impact_cost'].empty else 1000
                        fig_risk_matrix.update_layout(xaxis=dict(range=[0, 1]), yaxis=dict(title='Impact Cost ($)', range=[min_y, max_y]), template=PLOTLY_TEMPLATE, height=500)
                        fig_risk_matrix.add_vline(x=0.5, line_dash="dash", line_color=ARCADIS_GREY)
                        impact_75th = risk_matrix_df['impact_cost'].quantile(0.75) if not risk_matrix_df['impact_cost'].empty else 0
                        fig_risk_matrix.add_hline(y=impact_75th, line_dash="dash", line_color=ARCADIS_GREY, annotation_text="75th Percentile Impact", annotation_position="bottom right")
                        st.plotly_chart(fig_risk_matrix, use_container_width=True)
                        st.markdown("<p class='insight-text'><i class='fas fa-lightbulb'></i><b>Insight:</b> Risks in the top-right quadrant require urgent attention. Monitor risks near threshold lines.</p>", unsafe_allow_html=True)
                    except Exception as e: logging.error(f"Error creating risk matrix chart: {e}", exc_info=True); st.warning("⚠️ Could not generate Risk Matrix chart.")
                else: st.warning(f"⚠️ Risk data is empty or missing columns for Risk Matrix.")

                st.divider()
                if not risks_analysis.empty and all(col in risks_analysis.columns for col in required_cols_sunburst):
                    st.markdown("#### Risk Breakdown by Owner & Status")
                    try:
                        df_sunburst = risks_analysis.copy()
                        df_sunburst['risk_score'] = pd.to_numeric(df_sunburst['risk_score'], errors='coerce').fillna(0.0)
                        df_sunburst = df_sunburst[df_sunburst['risk_score'] > 0]
                        if not df_sunburst.empty:
                            fig_sunburst = px.sunburst(df_sunburst, path=['owner', 'risk_status'], values='risk_score', title='Risk Exposure ($) by Owner and Status', color='risk_status', color_discrete_map={'Open': COLOR_DANGER, 'Mitigating': COLOR_WARNING, 'Closed': COLOR_SUCCESS, 'Realized': COLOR_GREY, '(?)':'#DDDDDD'}, hover_data={'risk_score':':,'})
                            fig_sunburst.update_traces(textinfo='label+percent entry'); fig_sunburst.update_layout(height=500, margin = dict(t=50, l=25, r=25, b=25))
                            st.plotly_chart(fig_sunburst, use_container_width=True)
                            st.markdown("<p class='insight-text'><i class='fas fa-lightbulb'></i><b>Insight:</b> Identify risk concentration by owner and status. Large 'Open' or 'Mitigating' segments for specific owners might need support.</p>", unsafe_allow_html=True)
                        else: st.info("ℹ️ No risks with positive risk score found for Sunburst chart.")
                    except Exception as e: logging.error(f"Error creating risk sunburst chart: {e}", exc_info=True); st.warning("⚠️ Could not generate Risk Breakdown chart.")
                else: st.warning(f"⚠️ Risk data is empty or missing columns for Sunburst chart.")

                st.divider()
                st.markdown("#### Top 10 Open/Mitigating Risks by Score")
                if not risks_analysis.empty and 'risk_status' in risks_analysis.columns and 'risk_score' in risks_analysis.columns:
                    open_mitigating_risks = risks_analysis[risks_analysis['risk_status'].isin(['Open', 'Mitigating'])].copy()
                    open_mitigating_risks['risk_score'] = pd.to_numeric(open_mitigating_risks['risk_score'], errors='coerce').fillna(0)
                    if not open_mitigating_risks.empty:
                        if 'project_id' in projects_analysis.columns and 'project_name' in projects_analysis.columns and 'project_id' in open_mitigating_risks.columns:
                            try:
                                top_risks = open_mitigating_risks.nlargest(10, 'risk_score').merge(projects_analysis[['project_id', 'project_name']], on='project_id', how='left')
                                display_risk_cols = ['project_name', 'risk_description', 'probability', 'impact_cost', 'risk_score', 'risk_status', 'mitigation_plan', 'owner']
                                display_risk_cols_present = [col for col in display_risk_cols if col in top_risks.columns]
                                if display_risk_cols_present:
                                    df_display_risks = top_risks[display_risk_cols_present].rename(columns={'project_name': 'Project', 'risk_description': 'Risk', 'probability': 'P', 'impact_cost': 'Impact', 'risk_score': 'Score', 'risk_status':'Status', 'mitigation_plan':'Mitigation', 'owner':'Owner'})
                                    df_display_risks['P'] = pd.to_numeric(df_display_risks['P'], errors='coerce')
                                    df_display_risks['Impact'] = pd.to_numeric(df_display_risks['Impact'], errors='coerce')
                                    df_display_risks['Score'] = pd.to_numeric(df_display_risks['Score'], errors='coerce')
                                    st.dataframe(df_display_risks.style.format({'P': '{:.1%}', 'Impact': lambda x: format_currency(x, compact=True), 'Score': '{:,.0f}'}), use_container_width=True)
                                else: st.warning("⚠️ Could not display top risks table due to missing columns after merge.")
                            except Exception as e: logging.error(f"Error creating top risks table: {e}", exc_info=True); st.warning("⚠️ Could not generate Top Risks table.")
                        else: st.warning("⚠️ Cannot link risks to project names: 'project_id' or 'project_name' missing.")
                    else: st.info("ℹ️ No 'Open' or 'Mitigating' risks found in the filtered data.")
                else: st.warning("⚠️ Risk data is empty or missing 'risk_status'/'risk_score' columns.")


            # --- Resource Allocation ---
            elif analysis_options == "Resource Allocation":
                st.markdown("### Resource Assignment Overview")
                if not tasks_analysis.empty and 'resource' in tasks_analysis.columns:
                    try:
                        resource_counts = tasks_analysis['resource'].astype(str).value_counts().reset_index()
                        resource_counts.columns = ['resource', 'task_count']
                        fig_resource = px.bar(resource_counts, x='resource', y='task_count', title="Task Count Assigned by Resource/Team", text_auto=True, labels={'resource': 'Resource/Team', 'task_count': 'Number of Assigned Tasks'}, color_discrete_sequence=[ARCADIS_SECONDARY_PALETTE[0]]) # Teal
                        fig_resource.update_layout(template=PLOTLY_TEMPLATE, height=400)
                        st.plotly_chart(fig_resource, use_container_width=True)
                        st.markdown("<p class='insight-text'><i class='fas fa-lightbulb'></i><b>Insight:</b> Shows task distribution. High counts might indicate overallocation, but workload depends on effort/duration.</p>", unsafe_allow_html=True)
                        st.info("ℹ️ Detailed workload analysis requires effort/hour data per task, which is not currently available in the mock data.")
                    except Exception as e: logging.error(f"Error creating resource allocation chart: {e}", exc_info=True); st.warning("⚠️ Could not generate Resource Allocation chart.")
                else: st.warning("⚠️ No task data with 'resource' information available.")

            # --- Change Control ---
            elif analysis_options == "Change Control":
                st.markdown("### Change Request Analysis")
                required_cols_cr_pie = ['status']
                required_cols_cr_metrics = ['status', 'impact_cost', 'impact_schedule_days']
                required_cols_cr_table = ['project_id', 'description', 'status', 'impact_cost', 'impact_schedule_days', 'date_submitted']

                if not changes_analysis.empty:
                    col1, col2 = st.columns(2)
                    with col1:
                        if all(col in changes_analysis.columns for col in required_cols_cr_pie):
                            st.markdown(f"#### Status Distribution ({len(changes_analysis)} Total CRs)")
                            try:
                                status_counts = changes_analysis['status'].astype(str).value_counts()
                                if not status_counts.empty:
                                    fig_cr_status = px.pie(status_counts, values=status_counts.values, names=status_counts.index, title="CR Status", hole=0.4, color_discrete_sequence=px.colors.qualitative.Pastel)
                                    fig_cr_status.update_traces(textposition='inside', textinfo='percent+label'); fig_cr_status.update_layout(showlegend=True, margin=dict(t=50, b=20, l=20, r=20), height=350, legend_title_text='Status')
                                    st.plotly_chart(fig_cr_status, use_container_width=True)
                                else: st.info("ℹ️ No status info for CRs.")
                            except Exception as e: logging.error(f"Error creating CR status pie chart: {e}", exc_info=True); st.warning("⚠️ Could not generate CR Status chart.")
                        else: st.warning(f"⚠️ Missing columns for CR Status chart.")

                    with col2:
                         if all(col in changes_analysis.columns for col in required_cols_cr_metrics):
                            st.markdown("#### Aggregate Impact (Approved CRs)")
                            try:
                                approved_changes = changes_analysis[changes_analysis['status'] == 'Approved'].copy()
                                approved_changes['impact_cost'] = pd.to_numeric(approved_changes['impact_cost'], errors='coerce').fillna(0)
                                approved_changes['impact_schedule_days'] = pd.to_numeric(approved_changes['impact_schedule_days'], errors='coerce').fillna(0)
                                st.metric("Total Cost Impact", format_currency(approved_changes['impact_cost'].sum(), compact=True))
                                avg_sched_impact = approved_changes['impact_schedule_days'].mean() if not approved_changes.empty else 0
                                st.metric("Avg Schedule Impact", f"{avg_sched_impact:.1f} days")
                            except Exception as e: logging.error(f"Error calculating CR metrics: {e}", exc_info=True); st.warning("⚠️ Could not calculate CR metrics.")
                         else: st.warning(f"⚠️ Missing columns for CR Metrics.")

                    st.divider()
                    st.markdown("#### Change Request Log Details")
                    if 'project_id' in projects_analysis.columns and 'project_name' in projects_analysis.columns and 'project_id' in changes_analysis.columns and all(col in changes_analysis.columns for col in required_cols_cr_table if col != 'project_name'):
                        try:
                            changes_display = changes_analysis.merge(projects_analysis[['project_id', 'project_name']], on='project_id', how='left')
                            display_cr_cols = ['project_name', 'description', 'status', 'impact_cost', 'impact_schedule_days', 'date_submitted']
                            display_cr_cols_present = [col for col in display_cr_cols if col in changes_display.columns]
                            if display_cr_cols_present:
                                if 'date_submitted' in changes_display.columns: changes_display['date_submitted'] = pd.to_datetime(changes_display['date_submitted'], errors='coerce').dt.strftime('%Y-%m-%d')
                                if 'impact_cost' in changes_display.columns: changes_display['impact_cost'] = pd.to_numeric(changes_display['impact_cost'], errors='coerce')
                                if 'impact_schedule_days' in changes_display.columns: changes_display['impact_schedule_days'] = pd.to_numeric(changes_display['impact_schedule_days'], errors='coerce')
                                st.dataframe(changes_display[display_cr_cols_present].rename(columns={'project_name':'Project', 'description':'Description', 'status':'Status', 'impact_cost':'Cost Impact', 'impact_schedule_days':'Sched. Impact (d)', 'date_submitted':'Submitted'})
                                             .style.format({'Cost Impact': lambda x: format_currency(x, compact=True)}), use_container_width=True)
                                st.markdown("<p class='insight-text'><i class='fas fa-lightbulb'></i><b>Insight:</b> Monitor volume and impact of approved CRs. Analyze reasons for rejected/frequent changes.</p>", unsafe_allow_html=True)
                            else: st.warning("⚠️ Could not display CR table due to missing columns after merge.")
                        except Exception as e: logging.error(f"Error creating CR table: {e}", exc_info=True); st.warning("⚠️ Could not generate Change Request table.")
                    else: st.warning("⚠️ Cannot link changes to project names or required columns missing.")
                else: st.info("ℹ️ No change request data available in the filtered set.")

# --- Placeholder for Remaining Tabs (Implemented in Part 3) ---
# Add placeholders to avoid errors if Part 2 is run standalone
placeholder_tabs_p2 = ["🎯 Project Deep Dive", "💰 Benefits & ROI", "📝 Reports", "💾 Data & Settings"]
for tab_name in placeholder_tabs_p2:
     if f'{SESSION_PREFIX}tab_map_pmo' in st.session_state and tab_name in st.session_state[f'{SESSION_PREFIX}tab_map_pmo']:
         with st.session_state[f'{SESSION_PREFIX}tab_map_pmo'][tab_name]:
             st.header(tab_name)
             st.info(f"Content for {tab_name} will be implemented in Part 3.")

# Note: End of Part 2.
# -*- coding: utf-8 -*-
"""
PMO Pulse Narrative Driven App (v1.0 - Part 3)

Implements the Project Deep Dive, Benefits & ROI, Reports,
and Data & Settings tabs.
Relies on the structure and functions defined in Part 1 & 2.
"""

# -*- coding: utf-8 -*-
"""
PMO Pulse Narrative Driven App (v1.0 - Part 3 - KeyError Fix)

Implements the Project Deep Dive, Benefits & ROI, Reports,
and Data & Settings tabs.
Removed references to non-existent 'Variance_Pct' column in Reports tab.
Relies on the structure and functions defined in Part 1 & 2.
"""

# ==============================================================================
# Part 3: Tab Implementations (Project Deep Dive, Benefits, Reports, Data) & End
# ==============================================================================

import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import datetime
from dateutil.relativedelta import relativedelta
import time
import logging
import io
import os
# from fpdf import FPDF # Optional

# --- Assume Helper Functions & Constants are available from Part 1 ---
# (e.g., format_currency, create_gauge, SESSION_PREFIX, COLOR_*, PLOTLY_TEMPLATE)
# --- Assume Session State is initialized and populated from Part 1 & 2 ---

# --- Retrieve necessary variables from session state ---
# Use .get() with defaults for safety
SESSION_PREFIX = "pmo_pulse_" # Ensure prefix is defined
tab_map_pmo = st.session_state.get(f'{SESSION_PREFIX}tab_map_pmo', {})
data_loaded = st.session_state.get(f'{SESSION_PREFIX}data_loaded', False)
# Retrieve filtered data calculated in Part 1 sidebar logic
filtered_projects_with_kpis = st.session_state.get(f'{SESSION_PREFIX}filtered_projects_with_kpis', pd.DataFrame())
filtered_tasks = st.session_state.get(f'{SESSION_PREFIX}filtered_tasks', pd.DataFrame())
filtered_risks = st.session_state.get(f'{SESSION_PREFIX}filtered_risks', pd.DataFrame())
filtered_changes = st.session_state.get(f'{SESSION_PREFIX}filtered_changes', pd.DataFrame())
filtered_history = st.session_state.get(f'{SESSION_PREFIX}filtered_history', pd.DataFrame())
# Retrieve full dataframes as well (needed for some operations like full export)
all_dfs = st.session_state.get(f'{SESSION_PREFIX}all_dfs', {})
benefits_df = st.session_state.get(f'{SESSION_PREFIX}benefits_df', pd.DataFrame())

# --- Plotly Template (Define if not carried over) ---
ARCADIS_ORANGE = "#E67300"; ARCADIS_BLACK = "#000000"; ARCADIS_GREY = "#6c757d"; ARCADIS_WHITE = "#FFFFFF"; ARCADIS_DARK_GREY = "#646469"; ARCADIS_SECONDARY_PALETTE = ["#00A3A1", ARCADIS_DARK_GREY, "#D6D6D8"]; COLOR_SUCCESS = "#2ECC71"; COLOR_WARNING = "#F1C40F"; COLOR_DANGER = "#E74C3C"; COLOR_INFO = "#3498DB"; PLOTLY_TEMPLATE = "plotly_white"

# --- Formatting Helpers (Define if not carried over) ---
def format_currency(value, compact=False):
    if pd.isna(value) or not isinstance(value, (int, float, np.number)): return "$0"
    try:
        value = float(value); abs_value = abs(value); sign = "-" if value < 0 else ""
        if compact:
            if abs_value >= 1_000_000_000: return f"{sign}${abs_value / 1_000_000_000:.1f}B"
            if abs_value >= 1_000_000: return f"{sign}${abs_value / 1_000_000:.1f}M"
            if abs_value >= 1_000: return f"{sign}${abs_value / 1_000:.0f}K"
            return f"{sign}${abs_value:.0f}"
        else: return f"{sign}${value:,.0f}"
    except (ValueError, TypeError): return "$0"

def format_percentage(x):
    if pd.isna(x): return "N/A"
    try: return f"{x:.1f}%"
    except (ValueError, TypeError): return "Invalid"

# --- PowerPoint Reporting Functions (Assume from Part 1) ---
# Need access to these functions if defined in Part 1
# add_title_slide, add_kpi_slide, add_plot_slide, add_list_slide, generate_pptx_report

# Re-define generate_pptx_report to remove supplier analysis dependency
# (Or ensure the original definition handles missing supplier data gracefully)
from pptx import Presentation
from pptx.util import Inches, Pt

def add_title_slide(prs, title_text, subtitle_text):
    title_slide_layout = prs.slide_layouts[0]; slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title; subtitle = slide.placeholders[1]; title.text = title_text; subtitle.text = subtitle_text

def add_kpi_slide(prs, title_text, kpis):
    bullet_slide_layout = prs.slide_layouts[1]; slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes; title_shape = shapes.title; body_shape = shapes.placeholders[1]; title_shape.text = title_text
    tf = body_shape.text_frame; tf.clear(); tf.word_wrap = True
    for name, value, help_text in kpis:
        p = tf.add_paragraph(); run = p.add_run(); run.text = f"{name}: {value}"; run.font.bold = True; run.font.size = Pt(18)
        p_help = tf.add_paragraph(); run_help = p_help.add_run(); run_help.text = help_text; run_help.font.size = Pt(14); p_help.level = 1

def add_plot_slide(prs, title_text, plot_path, notes=""):
    pic_slide_layout = prs.slide_layouts[5]; slide = prs.slides.add_slide(pic_slide_layout)
    title_shape = slide.shapes.title; title_shape.text = title_text
    left, top, height = Inches(1), Inches(1.5), Inches(5.0)
    if plot_path and os.path.exists(plot_path):
         try: pic = slide.shapes.add_picture(plot_path, left, top, height=height)
         except Exception as e:
             left_txt, top_txt, width_txt, height_txt = Inches(1), Inches(1.8), Inches(8), Inches(1)
             txBox = slide.shapes.add_textbox(left_txt, top_txt, width_txt, height_txt); tf = txBox.text_frame
             tf.text = f"Error adding plot:\n{os.path.basename(plot_path)}\n{e}"
    else:
        left_txt, top_txt, width_txt, height_txt = Inches(1), Inches(1.8), Inches(8), Inches(1)
        txBox = slide.shapes.add_textbox(left_txt, top_txt, width_txt, height_txt); tf = txBox.text_frame
        tf.text = f"Plot image not available:\n'{title_text}'"
    if notes: notes_slide = slide.notes_slide; text_frame = notes_slide.notes_text_frame; text_frame.text = notes

def add_list_slide(prs, title_text, items_list, header=""):
    bullet_slide_layout = prs.slide_layouts[1]; slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes; title_shape = shapes.title; body_shape = shapes.placeholders[1]; title_shape.text = title_text
    tf = body_shape.text_frame; tf.clear(); tf.word_wrap = True
    if header: p = tf.add_paragraph(); run = p.add_run(); run.text = header; run.font.bold = True; run.font.size = Pt(16)
    for item in items_list: p = tf.add_paragraph(); run = p.add_run(); run.text = str(item); run.font.size = Pt(14); p.level = 1 if header else 0

# --- UPDATED generate_pptx_report ---
def generate_pptx_report(kpis, scatter_plot_path, variance_plot_path, outlier_plot_path): # Removed high_variance_suppliers
    """Generates the full PowerPoint report (without supplier analysis)."""
    prs = Presentation()
    add_title_slide(prs, "PMO Pulse Portfolio Report", f"Generated: {datetime.datetime.now().strftime('%Y-%m-%d')}\nArcadis Confidential")
    add_kpi_slide(prs, "Executive Summary - Key Performance Indicators", kpis)
    add_plot_slide(prs, "Project Performance Matrix (CPI vs SPI)", scatter_plot_path, notes="Shows overall performance. Size by budget.") # Updated title
    add_plot_slide(prs, "KPI Trends (SPI & CPI)", variance_plot_path, notes="Visualizes portfolio SPI and CPI trends over time.") # Updated title
    add_plot_slide(prs, "Top 5 Projects by Forecast Variance (VAC)", outlier_plot_path, notes="Highlights projects forecast to have the largest budget variance (negative = overrun).") # Updated title

    # Removed supplier slide

    conclusion_items = [
        "Review projects with low SPI/CPI or negative VAC for potential intervention.",
        "Analyze performance trends to identify systemic issues or improvements.",
        "Leverage detailed analysis tabs for further investigation into specific projects or segments.",
        "Contact the Arcadis Digital Team for further support or customization."
    ]
    add_list_slide(prs, "Conclusions & Next Steps", conclusion_items)
    pptx_io = io.BytesIO(); prs.save(pptx_io); pptx_io.seek(0); return pptx_io
# --- End of updated generate_pptx_report ---

@st.cache_data
def convert_df_to_csv(df):
    output = io.BytesIO(); df.to_csv(output, index=False, encoding='utf-8'); return output.getvalue()

@st.cache_data
def convert_df_to_excel(df_dict): # Modified to accept dict of dataframes
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        for sheet_name, df in df_dict.items():
            if isinstance(df, pd.DataFrame) and not df.empty:
                 # Convert datetime columns to timezone-naive if they exist
                 for col in df.select_dtypes(include=['datetime64[ns, UTC]', 'datetime64[ns]']).columns:
                     try: df[col] = df[col].dt.tz_localize(None)
                     except TypeError: pass # Already naive
                 df.to_excel(writer, index=False, sheet_name=sheet_name)
    return output.getvalue()


# --- Project Deep Dive Tab ---
project_tab_key = "🎯 Project Deep Dive"
if project_tab_key in tab_map_pmo:
    with tab_map_pmo[project_tab_key]:
        st.header("🎯 Project Deep Dive")
        st.markdown("_Select a specific project for detailed analysis of KPIs, tasks, risks, and changes._")

        if not data_loaded or filtered_projects_with_kpis.empty:
            st.info("ℹ️ No projects match the current filters or data is not loaded. Adjust filters in the sidebar or load data.")
        else:
            project_list_dd = sorted(filtered_projects_with_kpis['project_name'].unique().tolist())
            selected_project_name_dd = st.selectbox("Select Project:", project_list_dd, key="deep_dive_select_pmo")

            if selected_project_name_dd:
                try:
                    # Get project details (ensure only one row is selected)
                    selected_project_details_df = filtered_projects_with_kpis[filtered_projects_with_kpis['project_name'] == selected_project_name_dd]
                    if selected_project_details_df.empty:
                         st.error(f"Could not find details for project: {selected_project_name_dd}.")
                    else:
                        selected_project_details = selected_project_details_df.iloc[0]
                        selected_project_id = selected_project_details['project_id']

                        # Filter related data for the selected project
                        project_tasks_dd = filtered_tasks[filtered_tasks['project_id'] == selected_project_id] if 'project_id' in filtered_tasks else pd.DataFrame()
                        project_risks_dd = filtered_risks[filtered_risks['project_id'] == selected_project_id] if 'project_id' in filtered_risks else pd.DataFrame()
                        project_changes_dd = filtered_changes[filtered_changes['project_id'] == selected_project_id] if 'project_id' in filtered_changes else pd.DataFrame()
                        project_history_dd = filtered_history[filtered_history['project_id'] == selected_project_id] if 'project_id' in filtered_history else pd.DataFrame()

                        # --- Display Project Header ---
                        st.markdown(f"### {selected_project_details.get('project_name', 'N/A')}")
                        col1_ddh, col2_ddh, col3_ddh, col4_ddh = st.columns(4)
                        with col1_ddh: st.markdown(f"**Sector:** {selected_project_details.get('sector', 'N/A')}")
                        with col2_ddh: st.markdown(f"**PM:** {selected_project_details.get('project_manager', 'N/A')}")
                        with col3_ddh: st.markdown(f"**Status:** {selected_project_details.get('status', 'N/A')}")
                        with col4_ddh: st.markdown(f"**Alignment:** {selected_project_details.get('strategic_alignment', 'N/A')}")

                        if 'collaboration_link' in selected_project_details and pd.notna(selected_project_details['collaboration_link']):
                             st.markdown(f"**Link:** [Project Details]({selected_project_details['collaboration_link']})", unsafe_allow_html=True)

                        planned_end_date_str = pd.to_datetime(selected_project_details.get('planned_end_date', pd.NaT)).strftime('%Y-%m-%d') if pd.notna(selected_project_details.get('planned_end_date', pd.NaT)) else 'N/A'
                        st.markdown(f"**Budget:** {format_currency(selected_project_details.get('budget', 0))} | **Planned End:** {planned_end_date_str}")
                        st.divider()

                        # --- Display KPIs ---
                        st.markdown("#### Key Performance Indicators")
                        kpi_cols_dd = st.columns(4)
                        with kpi_cols_dd[0]: st.metric("SPI", f"{selected_project_details.get('spi', 0):.2f}")
                        with kpi_cols_dd[1]: st.metric("CPI", f"{selected_project_details.get('cpi', 0):.2f}")
                        with kpi_cols_dd[2]: st.metric("Cost Variance (CV)", format_currency(selected_project_details.get('cv', 0)))
                        with kpi_cols_dd[3]: st.metric("Schedule Variance (SV)", format_currency(selected_project_details.get('sv', 0)))

                        fin_cols_dd = st.columns(3)
                        with fin_cols_dd[0]: st.metric("Budget (BAC)", format_currency(selected_project_details.get('bac', 0)))
                        with fin_cols_dd[1]: st.metric("Forecast (EAC)", format_currency(selected_project_details.get('eac_cpi', 0)))
                        with fin_cols_dd[2]: st.metric("Variance (VAC)", format_currency(selected_project_details.get('vac', 0), compact=True))

                        # --- Display Performance Trend ---
                        if not project_history_dd.empty and 'month' in project_history_dd.columns and 'spi' in project_history_dd.columns and 'cpi' in project_history_dd.columns:
                            st.markdown("#### Performance Trend (SPI/CPI)")
                            try:
                                project_history_numeric = project_history_dd.copy()
                                project_history_numeric['spi'] = pd.to_numeric(project_history_numeric['spi'], errors='coerce')
                                project_history_numeric['cpi'] = pd.to_numeric(project_history_numeric['cpi'], errors='coerce')
                                project_history_numeric = project_history_numeric.dropna(subset=['spi', 'cpi', 'month'])
                                if not project_history_numeric.empty:
                                    project_history_sorted = project_history_numeric.sort_values('month')
                                    fig_proj_trend = go.Figure()
                                    fig_proj_trend.add_trace(go.Scatter(x=project_history_sorted['month'], y=project_history_sorted['spi'], mode='lines+markers', name='SPI', line=dict(color=ARCADIS_SECONDARY_PALETTE[0]))) # Teal
                                    fig_proj_trend.add_trace(go.Scatter(x=project_history_sorted['month'], y=project_history_sorted['cpi'], mode='lines+markers', name='CPI', line=dict(color=ARCADIS_ORANGE)))
                                    fig_proj_trend.add_hline(y=1.0, line_dash="dash", line_color=ARCADIS_GREY)
                                    fig_proj_trend.update_layout(title=None, xaxis_title=None, yaxis_title="Index", height=300, margin=dict(t=20, b=30, l=30, r=30), template=PLOTLY_TEMPLATE, legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1))
                                    st.plotly_chart(fig_proj_trend, use_container_width=True)
                                else: st.info("ℹ️ No valid historical KPI data available for this project after cleaning.")
                            except Exception as e: logging.error(f"Error creating project trend chart for {selected_project_name_dd}: {e}", exc_info=True); st.warning("⚠️ Could not generate project performance trend chart.")
                        else: st.info("ℹ️ No historical KPI data available for this project.")

                        # --- Display Tasks, Risks, Changes in Expanders ---
                        st.divider()
                        st.markdown("#### Project Details")
                        exp_cols_dd = st.columns(3)
                        with exp_cols_dd[0]:
                            with st.expander(f"Tasks ({len(project_tasks_dd)})"):
                                if not project_tasks_dd.empty:
                                    display_task_cols = ['task_name', 'planned_start', 'planned_end', 'actual_start', 'actual_end', 'percent_complete', 'planned_cost', 'actual_cost', 'earned_value', 'resource']
                                    display_task_cols_present = [col for col in display_task_cols if col in project_tasks_dd.columns]
                                    tasks_display_df = project_tasks_dd[display_task_cols_present].copy()
                                    date_cols = ['planned_start', 'planned_end', 'actual_start', 'actual_end']
                                    for col in date_cols:
                                        if col in tasks_display_df.columns: tasks_display_df[col] = pd.to_datetime(tasks_display_df[col], errors='coerce').dt.strftime('%Y-%m-%d')
                                    if 'percent_complete' in tasks_display_df.columns:
                                        tasks_display_df['percent_complete'] = pd.to_numeric(tasks_display_df['percent_complete'], errors='coerce')
                                        tasks_display_df['percent_complete'] = tasks_display_df['percent_complete'].apply(lambda x: f"{x:.0%}" if pd.notna(x) else 'N/A')
                                    tasks_display_df = tasks_display_df.rename(columns={'task_name':'Task', 'planned_start':'Plan Start', 'planned_end':'Plan End', 'actual_start':'Actual Start', 'actual_end':'Actual End', 'percent_complete':'% Comp', 'planned_cost':'Plan Cost', 'actual_cost':'Actual Cost', 'earned_value':'EV', 'resource':'Resource'})
                                    column_config_tasks = {"Plan Cost": st.column_config.NumberColumn(format="$ ,.0f"), "Actual Cost": st.column_config.NumberColumn(format="$ ,.0f"), "EV": st.column_config.NumberColumn(format="$ ,.0f")}
                                    st.dataframe(tasks_display_df, column_config=column_config_tasks, height=300, use_container_width=True)
                                else: st.info("ℹ️ No tasks found.")
                        with exp_cols_dd[1]:
                            with st.expander(f"Risks ({len(project_risks_dd)})"):
                                if not project_risks_dd.empty:
                                    display_risk_cols = ['risk_description', 'probability', 'impact_cost', 'risk_score', 'risk_status', 'mitigation_plan', 'owner', 'collaboration_link']
                                    display_risk_cols_present = [col for col in display_risk_cols if col in project_risks_dd.columns]
                                    df_display_risks = project_risks_dd[display_risk_cols_present].rename(columns={'risk_description':'Risk', 'probability':'P', 'impact_cost':'Impact', 'risk_score':'Score', 'risk_status':'Status', 'mitigation_plan':'Mitigation', 'collaboration_link':'Link'})
                                    df_display_risks['P'] = pd.to_numeric(df_display_risks['P'], errors='coerce')
                                    df_display_risks['Impact'] = pd.to_numeric(df_display_risks['Impact'], errors='coerce')
                                    df_display_risks['Score'] = pd.to_numeric(df_display_risks['Score'], errors='coerce')
                                    column_config_risks = {"P": st.column_config.NumberColumn(format="%.1f%%"), "Impact": st.column_config.NumberColumn(format="$ ,.0f", help="Impact Cost"), "Score": st.column_config.NumberColumn(format=",.0f")}
                                    if 'Link' in df_display_risks.columns: column_config_risks["Link"] = st.column_config.LinkColumn("Link", display_text="Details")
                                    st.dataframe(df_display_risks, column_config=column_config_risks, height=300, use_container_width=True)
                                else: st.info("ℹ️ No risks found.")
                        with exp_cols_dd[2]:
                            with st.expander(f"Change Requests ({len(project_changes_dd)})"):
                                if not project_changes_dd.empty:
                                    display_change_cols = ['description', 'status', 'impact_cost', 'impact_schedule_days', 'date_submitted', 'collaboration_link']
                                    display_change_cols_present = [col for col in display_change_cols if col in project_changes_dd.columns]
                                    df_display_changes = project_changes_dd[display_change_cols_present].rename(columns={'description':'Description','impact_cost':'Cost Impact', 'impact_schedule_days':'Sched. Impact (d)', 'date_submitted':'Submitted', 'collaboration_link':'Link'})
                                    if 'Submitted' in df_display_changes.columns: df_display_changes['Submitted'] = pd.to_datetime(df_display_changes['Submitted'], errors='coerce').dt.strftime('%Y-%m-%d')
                                    if 'Cost Impact' in df_display_changes.columns: df_display_changes['Cost Impact'] = pd.to_numeric(df_display_changes['Cost Impact'], errors='coerce')
                                    if 'Sched. Impact (d)' in df_display_changes.columns: df_display_changes['Sched. Impact (d)'] = pd.to_numeric(df_display_changes['Sched. Impact (d)'], errors='coerce')
                                    column_config_changes = {"Cost Impact": st.column_config.NumberColumn(format="$ ,.0f"), "Sched. Impact (d)": st.column_config.NumberColumn(format="%d d")}
                                    if 'Link' in df_display_changes.columns: column_config_changes["Link"] = st.column_config.LinkColumn("Link", display_text="Details")
                                    st.dataframe(df_display_changes, column_config=column_config_changes, height=300, use_container_width=True)
                                else: st.info("ℹ️ No change requests found.")

                except IndexError:
                    st.error(f"Could not retrieve details for project: {selected_project_name_dd}. Please try re-selecting.")
                    logging.warning(f"IndexError retrieving details for project: {selected_project_name_dd}")
                except Exception as e:
                    logging.error(f"Error rendering deep dive for {selected_project_name_dd}: {e}", exc_info=True)
                    st.error(f"An error occurred while displaying details for {selected_project_name_dd}.")

            else:
                st.info("ℹ️ Select a project from the list above.")


# --- Benefits & ROI Tab ---
benefits_tab_key = "💰 Benefits & ROI"
if benefits_tab_key in tab_map_pmo:
    with tab_map_pmo[benefits_tab_key]:
        st.header("💰 Benefits Realization & ROI")
        st.markdown("_Tracking the value delivered by improved PMO processes and this tool (using simulated data)._")

        required_cols_benefits = ['Month', 'ReportingTimeSaved_hrs', 'CostOverrunsAvoided_k', 'ForecastAccuracy_perc']

        if not benefits_df.empty and all(col in benefits_df.columns for col in required_cols_benefits):
            try:
                benefits_df_processed = benefits_df.copy()
                benefits_df_processed['Month_dt'] = pd.to_datetime(benefits_df_processed['Month'] + '-01', errors='coerce')
                benefits_df_processed['ReportingTimeSaved_hrs'] = pd.to_numeric(benefits_df_processed['ReportingTimeSaved_hrs'], errors='coerce').fillna(0)
                benefits_df_processed['CostOverrunsAvoided_k'] = pd.to_numeric(benefits_df_processed['CostOverrunsAvoided_k'], errors='coerce').fillna(0)
                benefits_df_processed['ForecastAccuracy_perc'] = pd.to_numeric(benefits_df_processed['ForecastAccuracy_perc'], errors='coerce').fillna(0)
                benefits_df_processed = benefits_df_processed.dropna(subset=['Month_dt'])

                if benefits_df_processed.empty:
                     st.warning("⚠️ Benefits data is empty after processing.")
                else:
                    latest_benefits = benefits_df_processed.iloc[-1]
                    with st.container(): # Use container for layout
                        col1, col2, col3 = st.columns(3)
                        with col1: st.metric("Reporting Time Saved (Last Month)", f"{latest_benefits.get('ReportingTimeSaved_hrs', 0):.1f} hrs")
                        with col2: st.metric("Cumulative Cost Avoidance", f"{format_currency(latest_benefits.get('CostOverrunsAvoided_k', 0) * 1000)}")
                        with col3: st.metric("Forecast Accuracy (Last Month)", f"{latest_benefits.get('ForecastAccuracy_perc', 0):.1f}%")

                    st.divider()
                    st.markdown("### Benefit Trends (Simulated)")
                    chart_cols = st.columns(3)
                    with chart_cols[0]:
                        fig_ben1 = px.line(benefits_df_processed, x='Month_dt', y='ReportingTimeSaved_hrs', title='Reporting Time Saved', markers=True, color_discrete_sequence=[COLOR_SECONDARY])
                        fig_ben1.update_layout(xaxis_title=None, yaxis_title="Hours Saved", height=280, margin=dict(t=40, b=20, l=40, r=20), template=PLOTLY_TEMPLATE, title_font_size=16)
                        st.plotly_chart(fig_ben1, use_container_width=True)
                    with chart_cols[1]:
                        fig_ben2 = px.line(benefits_df_processed, x='Month_dt', y='CostOverrunsAvoided_k', title='Cost Avoidance ($K)', markers=True, color_discrete_sequence=[COLOR_SUCCESS])
                        fig_ben2.update_layout(xaxis_title=None, yaxis_title="Cumulative $K Avoided", height=280, margin=dict(t=40, b=20, l=40, r=20), template=PLOTLY_TEMPLATE, title_font_size=16)
                        st.plotly_chart(fig_ben2, use_container_width=True)
                    with chart_cols[2]:
                        fig_ben3 = px.line(benefits_df_processed, x='Month_dt', y='ForecastAccuracy_perc', title='Forecast Accuracy (%)', markers=True, color_discrete_sequence=[COLOR_INFO])
                        min_acc = max(0, benefits_df_processed['ForecastAccuracy_perc'].min() - 5)
                        max_acc = min(100, benefits_df_processed['ForecastAccuracy_perc'].max() + 5)
                        fig_ben3.update_layout(xaxis_title=None, yaxis_title="Accuracy %", yaxis_range=[min_acc, max_acc], height=280, margin=dict(t=40, b=20, l=40, r=20), template=PLOTLY_TEMPLATE, title_font_size=16)
                        st.plotly_chart(fig_ben3, use_container_width=True)

                    st.divider()
                    st.subheader("Return on Investment (ROI) - Placeholder")
                    st.info("ℹ️ Calculating ROI requires defining investment cost and quantifying financial value of benefits.")
                    roi_cols = st.columns([1, 2])
                    with roi_cols[0]:
                        invest_cost = st.number_input("Estimated Investment Cost ($)", min_value=0, value=50000, step=10000, key="roi_input")
                    with roi_cols[1]:
                        total_financial_benefit = latest_benefits.get('CostOverrunsAvoided_k', 0) * 1000
                        if invest_cost > 0:
                            roi_perc = ((total_financial_benefit - invest_cost) / invest_cost) * 100
                            st.metric("Simple ROI (Cost Avoidance vs Investment)", f"{roi_perc:.1f}%")
                        else:
                            st.metric("Simple ROI (Cost Avoidance vs Investment)", "N/A", help="Enter Investment Cost > 0")
            except Exception as e:
                logging.error(f"Error rendering benefits tab: {e}", exc_info=True)
                st.error("An error occurred while displaying benefits information.")
        else:
            missing_cols = [col for col in required_cols_benefits if col not in benefits_df.columns] if isinstance(benefits_df, pd.DataFrame) else required_cols_benefits
            st.warning(f"⚠️ Benefits tracking data is not available or missing required columns: {missing_cols}")

# --- Reports Tab ---
reports_tab_key = "📝 Reports"
if reports_tab_key in tab_map_pmo:
    with tab_map_pmo[reports_tab_key]:
        st.header("📝 Reports")
        st.markdown("_Generate summary reports for the currently filtered portfolio._")

        if not data_loaded or filtered_projects_with_kpis.empty:
            st.info("ℹ️ No projects match the current filters or data is not loaded. Cannot generate reports.")
        else:
            df_report_data_rep = filtered_projects_with_kpis.copy() # Use filtered data
            st.markdown(f"Report will be based on the **{len(df_report_data_rep)}** currently filtered projects.")

            # --- Generate PowerPoint Report ---
            st.subheader("Generate Executive PowerPoint Report")
            st.markdown("Create a downloadable PPTX summary of the key findings for the **filtered** data.")

            # Prepare KPIs for the filtered report
            report_kpis_rep = []
            if not df_report_data_rep.empty:
                 # Use .get() for KPIs in case calculation failed
                 avg_spi_rep = df_report_data_rep['spi'].mean()
                 avg_cpi_rep = df_report_data_rep['cpi'].mean()
                 num_over_rep = len(df_report_data_rep[pd.to_numeric(df_report_data_rep['vac'], errors='coerce') < 0])
                 total_vac_rep = df_report_data_rep['vac'].sum() # Sum of VAC across filtered projects

                 report_kpis_rep = [
                    ("Avg. Schedule Performance (SPI)", f"{avg_spi_rep:.2f}", "Average SPI for filtered projects."),
                    ("Avg. Cost Performance (CPI)", f"{avg_cpi_rep:.2f}", "Average CPI for filtered projects."),
                    ("# Projects Forecast Over Budget", f"{num_over_rep} / {len(df_report_data_rep)}", "Count of projects with negative VAC (Filtered)."),
                    ("Total Portfolio VAC ($M)", format_currency(total_vac_rep), "Sum of Variance at Completion for filtered projects."),
                 ]

            # Retrieve plot paths from session state (saved in Exec Summary)
            # These plots show OVERALL portfolio, not filtered. Consider regenerating if filtered view is needed.
            scatter_path_rep = st.session_state.get('scatter_plot_path')
            variance_path_rep = st.session_state.get('variance_plot_path') # This was the KPI trend plot
            outlier_plot_path_rep = st.session_state.get('outlier_plot_path') # This was the VAC plot

            scatter_plot_exists_rep = scatter_path_rep and os.path.exists(scatter_path_rep)
            variance_plot_exists_rep = variance_path_rep and os.path.exists(variance_path_rep)
            outlier_plot_exists_rep = outlier_plot_path_rep and os.path.exists(outlier_plot_path_rep)
            plots_available_rep = all([scatter_plot_exists_rep, variance_plot_exists_rep, outlier_plot_exists_rep])

            if not plots_available_rep:
                 st.warning("Some summary plots (based on unfiltered data) needed for the report were not generated or saved correctly. Report may lack images.")

            if st.button("Generate Filtered Report (.pptx)", key="gen_pptx_rep"):
                if not report_kpis_rep:
                     st.error("Cannot generate report: KPIs could not be calculated for filtered data.")
                else:
                    with st.spinner("Generating PowerPoint report..."):
                         try:
                             # Pass relevant info to the updated report function
                             pptx_file_rep = generate_pptx_report(
                                 kpis=report_kpis_rep,
                                 scatter_plot_path=scatter_path_rep if scatter_plot_exists_rep else None,
                                 variance_plot_path=variance_path_rep if variance_plot_exists_rep else None, # KPI Trend plot path
                                 outlier_plot_path=outlier_plot_path_rep if outlier_plot_exists_rep else None # VAC plot path
                                 # high_variance_suppliers removed
                             )
                             st.download_button("📥 Download Filtered Report (.pptx)", pptx_file_rep, f"Arcadis_PMO_Pulse_Filtered_Report_{datetime.datetime.now().strftime('%Y%m%d')}.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", key="dl_pptx_rep")
                         except Exception as e:
                             st.error(f"Failed to generate report: {e}")
                             import traceback; st.error(traceback.format_exc())

# --- Data & Settings Tab ---
data_settings_tab_key = "💾 Data & Settings"
if data_settings_tab_key in tab_map_pmo:
    with tab_map_pmo[data_settings_tab_key]:
        st.header("💾 Data & Settings")
        st.markdown("_Manage data sources, exports, and view application settings._")

        # --- Data Upload/Management ---
        st.subheader("📤 Upload & Manage Data")
        st.info("Use the sidebar controls for quick data refresh or mock data reset. Use this section for detailed upload or template download.")
        with st.container():
            data_type_ds = st.selectbox("Data Type to Upload:", ["Projects", "Tasks", "Risks", "Changes", "History", "Benefits"], key="data_type_upload_ds")
            uploaded_file_ds = st.file_uploader(f"Upload {data_type_ds} (CSV/Excel)", type=['csv', 'xlsx'], key="data_uploader_ds")

            if uploaded_file_ds:
                with st.spinner(f'Processing {data_type_ds}...'):
                    try:
                        if uploaded_file_ds.name.endswith('.csv'): df_uploaded_ds = pd.read_csv(uploaded_file_ds)
                        else: df_uploaded_ds = pd.read_excel(uploaded_file_ds)

                        # Basic validation based on type (replace with more robust checks)
                        df_key = data_type_ds.lower()
                        required_cols_ds = { # Define required cols for validation
                            'projects': ['project_id', 'project_name', 'sector', 'budget', 'start_date', 'planned_end_date', 'status'],
                            'tasks': ['task_id', 'project_id', 'planned_start', 'planned_end', 'planned_cost', 'actual_cost', 'earned_value'],
                            'risks': ['risk_id', 'project_id', 'probability', 'impact_cost'],
                            'changes': ['change_id', 'project_id', 'status', 'impact_cost', 'impact_schedule_days'],
                            'history': ['project_id', 'month', 'cpi', 'spi'],
                            'benefits': ['Month', 'ReportingTimeSaved_hrs', 'CostOverrunsAvoided_k', 'ForecastAccuracy_perc']
                        }.get(df_key, [])

                        is_valid_ds, errors_ds, warnings_ds, validated_data_ds = validate_data(df_uploaded_ds, df_key) # Pass required cols if needed
                        [st.warning(w) for w in warnings_ds]

                        if is_valid_ds:
                            st.session_state[f'{SESSION_PREFIX}all_dfs'][df_key] = validated_data_ds # Update the specific dataframe in the dict
                            st.session_state[f'{SESSION_PREFIX}{df_key}_df'] = validated_data_ds # Update the specific session state variable
                            st.success(f"{data_type_ds} data uploaded successfully. Recalculating KPIs...")
                            # Recalculate KPIs and rerun
                            st.cache_data.clear() # Clear cache as data changed
                            st.session_state[f'{SESSION_PREFIX}data_loaded'] = True # Mark as loaded
                            st.rerun()
                        else:
                            [st.error(e) for e in errors_ds]

                    except Exception as e: st.error(f"Error processing file: {e}")

        st.markdown("---")
        # --- Data Preview ---
        st.subheader("📊 Data Preview")
        with st.container():
            preview_type_ds = st.selectbox("Preview Data:", list(st.session_state[f'{SESSION_PREFIX}all_dfs'].keys()), key="data_preview_select_ds")
            df_to_preview = st.session_state[f'{SESSION_PREFIX}all_dfs'].get(preview_type_ds)

            if df_to_preview is not None and not df_to_preview.empty:
                st.dataframe(df_to_preview.head(20), use_container_width=True) # Show more rows
                st.markdown(f"**Total Rows:** {len(df_to_preview)}")
            else: st.info(f"No data available for '{preview_type_ds}'.")

        st.markdown("---")
        # --- Download Templates ---
        st.subheader("📥 Download Data Templates")
        with st.container():
            st.markdown("Download sample CSV templates for uploading your own data.")
            template_cols_ds = {
                'Projects': ['project_id', 'project_name', 'sector', 'budget', 'start_date', 'planned_end_date', 'status', 'project_manager', 'planned_duration_days', 'target_cpi', 'target_spi', 'strategic_alignment', 'collaboration_link'],
                'Tasks': ['task_id', 'project_id', 'task_name', 'planned_start', 'planned_end', 'actual_start', 'actual_end', 'planned_cost', 'actual_cost', 'earned_value', 'percent_complete', 'resource'],
                'Risks': ['risk_id', 'project_id', 'risk_description', 'probability', 'impact_cost', 'risk_score', 'mitigation_plan', 'risk_status', 'owner', 'collaboration_link'],
                'Changes': ['change_id', 'project_id', 'description', 'impact_cost', 'impact_schedule_days', 'status', 'date_submitted', 'collaboration_link'],
                'History': ['project_id', 'month', 'cpi', 'spi'],
                'Benefits': ['Month', 'ReportingTimeSaved_hrs', 'CostOverrunsAvoided_k', 'ForecastAccuracy_perc']
            }
            for data_type, cols in template_cols_ds.items():
                template_df_ds = pd.DataFrame(columns=cols)
                csv_bytes_ds = convert_df_to_csv(template_df_ds)
                st.download_button(
                    label=f"Download {data_type} Template", data=csv_bytes_ds,
                    file_name=f"{data_type.lower()}_template.csv", mime='text/csv',
                    key=f"download_template_{data_type.lower()}_ds"
                )

        st.markdown("---")
        # --- Export Full Dataset ---
        st.subheader("📤 Export Full Dataset")
        st.markdown("Export all currently loaded data (unfiltered) to a single Excel file.")
        if st.button("Export All Data (Excel)", key="export_all_excel"):
             with st.spinner("Preparing full Excel export..."):
                all_dfs_export = st.session_state.get(f'{SESSION_PREFIX}all_dfs', {})
                if all_dfs_export:
                    try:
                        # Pass the dictionary of dataframes to the export function
                        excel_data_full = convert_df_to_excel(all_dfs_export)
                        st.download_button(
                            label="📥 Download Full Export (.xlsx)",
                            data=excel_data_full,
                            file_name=f"pmo_pulse_full_export_{datetime.datetime.now().strftime('%Y%m%d')}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            key='dl_excel_full'
                        )
                    except Exception as e:
                         logging.error(f"Error creating full Excel export: {e}", exc_info=True)
                         st.error("⚠️ Could not generate full Excel export file.")
                else:
                    st.warning("No data loaded to export.")

# --- End of App ---
# Footer is now defined in Part 1
# st.markdown("---")
# st.caption("_PMO Pulse - Narrative Driven Analysis for Arcadis_")

# --- Main Execution Block (Optional - for direct script run) ---
# if __name__ == "__main__":
#     # Basic check for dependencies
#     try:
#         import streamlit; import pandas; import numpy; import plotly; import dateutil; import io
#         logging.info("Core libraries imported successfully.")
#     except ImportError as e:
#          logging.critical(f"Missing critical library: {e}. Please install required packages.")
#          st.error(f"Missing critical library: {e}. Please install required packages.")
#          st.stop()
#
#     # Authentication could be called here if needed before main()
#     # if authenticate_user():
#     #     main() # Assume main() is defined elsewhere or encompasses the tab logic
#     # else:
#     #     st.error("Authentication failed.")

# Note: The main() function structure from the original script is removed as
# Streamlit executes the script top-down. The tab rendering logic is now
# directly under the tab definitions.

